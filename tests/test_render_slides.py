"""Tests for ingestion.render_slides (B1)."""

from __future__ import annotations

import io
import shutil
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation

from ingestion.render_slides import (
    RenderSlidesError,
    _resolve_soffice,
    render_slides,
)


def _mini_pptx(slide_count: int = 3) -> bytes:
    prs = Presentation()
    # Default template has layouts; add blank-ish slides via layout 6 if present.
    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    # Presentation() starts with 0 slides in python-pptx.
    for _ in range(slide_count):
        prs.slides.add_slide(layout)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_resolve_soffice_missing(monkeypatch):
    monkeypatch.delenv("SOFFICE_BIN", raising=False)
    monkeypatch.setattr("ingestion.render_slides.shutil.which", lambda _n: None)
    real_is_file = Path.is_file

    def fake_is_file(self: Path) -> bool:
        if "LibreOffice.app" in str(self):
            return False
        return real_is_file(self)

    monkeypatch.setattr(Path, "is_file", fake_is_file)
    with pytest.raises(RenderSlidesError, match="LibreOffice not found"):
        _resolve_soffice()


def test_resolve_soffice_env(monkeypatch, tmp_path):
    bin_path = tmp_path / "soffice"
    bin_path.write_text("#!/bin/sh\n")
    monkeypatch.setenv("SOFFICE_BIN", str(bin_path))
    assert _resolve_soffice() == str(bin_path)


def test_render_slides_empty_indices(tmp_path):
    pptx = tmp_path / "deck.pptx"
    pptx.write_bytes(_mini_pptx(2))
    with (
        patch("ingestion.render_slides._resolve_soffice", return_value="/bin/soffice"),
        patch("ingestion.render_slides._resolve_pdftoppm", return_value="/bin/pdftoppm"),
        pytest.raises(ValueError, match="empty"),
    ):
        render_slides(pptx, [])


def test_render_slides_out_of_range(tmp_path):
    pptx = tmp_path / "deck.pptx"
    pptx.write_bytes(_mini_pptx(2))
    with (
        patch("ingestion.render_slides._resolve_soffice", return_value="/bin/soffice"),
        patch("ingestion.render_slides._resolve_pdftoppm", return_value="/bin/pdftoppm"),
        pytest.raises(ValueError, match="out of range"),
    ):
        render_slides(pptx, [0, 5])


def test_render_slides_missing_pptx(tmp_path):
    with (
        patch("ingestion.render_slides._resolve_soffice", return_value="/bin/soffice"),
        patch("ingestion.render_slides._resolve_pdftoppm", return_value="/bin/pdftoppm"),
        pytest.raises(FileNotFoundError),
    ):
        render_slides(tmp_path / "nope.pptx", [0])


def test_render_slides_happy_path(tmp_path, monkeypatch):
    pptx_bytes = _mini_pptx(3)
    out_dir = tmp_path / "pngs"

    def fake_convert(pptx_path, out_dir_work, *, soffice_bin, timeout):
        pdf = out_dir_work / f"{pptx_path.stem}.pdf"
        pdf.write_bytes(b"%PDF-1.4 fake")
        return pdf

    def fake_pdf_page(pdf_path, page_1based, dest_png, *, pdftoppm, dpi, timeout):
        # PNG magic + minimal payload
        dest_png.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 16)

    monkeypatch.setattr(
        "ingestion.render_slides._resolve_soffice", lambda _e=None: "/bin/soffice"
    )
    monkeypatch.setattr(
        "ingestion.render_slides._resolve_pdftoppm", lambda: "/bin/pdftoppm"
    )
    monkeypatch.setattr("ingestion.render_slides._convert_to_pdf", fake_convert)
    monkeypatch.setattr("ingestion.render_slides._pdf_page_to_png", fake_pdf_page)

    paths = render_slides(pptx_bytes, [0, 2], output_dir=out_dir)
    assert [p.name for p in paths] == ["slide-000.png", "slide-002.png"]
    assert all(p.is_file() for p in paths)
    assert paths[0].read_bytes()[:8] == b"\x89PNG\r\n\x1a\n"


def test_render_slides_all_indices_when_none(tmp_path, monkeypatch):
    pptx = tmp_path / "deck.pptx"
    pptx.write_bytes(_mini_pptx(2))
    seen: list[int] = []

    def fake_convert(pptx_path, out_dir_work, *, soffice_bin, timeout):
        pdf = out_dir_work / "deck.pdf"
        pdf.write_bytes(b"%PDF")
        return pdf

    def fake_pdf_page(pdf_path, page_1based, dest_png, *, pdftoppm, dpi, timeout):
        seen.append(page_1based)
        dest_png.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 8)

    monkeypatch.setattr(
        "ingestion.render_slides._resolve_soffice", lambda _e=None: "/bin/soffice"
    )
    monkeypatch.setattr(
        "ingestion.render_slides._resolve_pdftoppm", lambda: "/bin/pdftoppm"
    )
    monkeypatch.setattr("ingestion.render_slides._convert_to_pdf", fake_convert)
    monkeypatch.setattr("ingestion.render_slides._pdf_page_to_png", fake_pdf_page)

    paths = render_slides(pptx, None, output_dir=tmp_path / "out")
    assert seen == [1, 2]
    assert [p.name for p in paths] == ["slide-000.png", "slide-001.png"]


@pytest.mark.skipif(
    shutil.which("soffice") is None and shutil.which("libreoffice") is None,
    reason="LibreOffice not installed on host",
)
@pytest.mark.skipif(shutil.which("pdftoppm") is None, reason="pdftoppm not installed")
def test_render_slides_integration(tmp_path):
    pptx = tmp_path / "deck.pptx"
    pptx.write_bytes(_mini_pptx(2))
    out = tmp_path / "out"
    paths = render_slides(pptx, [1], output_dir=out, timeout=120)
    assert len(paths) == 1
    assert paths[0].name == "slide-001.png"
    assert paths[0].read_bytes()[:8] == b"\x89PNG\r\n\x1a\n"
