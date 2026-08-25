"""Integration smoke: Prodie-selected products → confirm_mix → build_deck.

  PYTHONPATH=. uv run python tests/smoke_ideation_to_deck.py --mock-ai

Uses a representative product selection as if Prodie + the associate checkboxes
already chose it. build_deck assembly uses real S3 FortuneAI + GTM when available;
--mock-ai skips Claude.
"""

from __future__ import annotations

import argparse
import io
import os
import re
import sys

import requests
from dotenv import load_dotenv
from pptx import Presentation

from ingestion.confirm_mix import confirm_mix_from_dict
from ingestion.generator import DeckGenerator
from ingestion.schema import DeckSchema
from tests.fortuneai_placeholder_fixture import MINIMAL_PNG, mock_placeholder_ai
from tests.logic_guide_fixtures import base_discovery_fields, build_representative_engine
from tests.smoke_build_live import LEFTOVER_TOKENS, _slide_text

# Prodie-selected newsletter with a budget aligned to the authoritative fixture price.
_DISCOVERY = base_discovery_fields(
    preferred_platforms_products=["Newsletters"],
    budgets=[{"amount": 5_000.0, "label": "Primary"}],
    targeting_details="CEO and C-suite leadership audience",
)


def _align_budgets(deck_schema: dict, mix_total: float) -> dict:
    """Set stated budget to mix total so C2 investment fill passes."""
    deck_schema = dict(deck_schema)
    deck_schema["budgets"] = [{"amount": mix_total, "label": "Primary"}]
    return deck_schema


def _run_selection_chain() -> tuple[dict, dict]:
    engine = build_representative_engine()
    discovery = _DISCOVERY
    gtm = engine._gtm
    inventory = engine._inventory
    confirm = confirm_mix_from_dict(
        {
            "discovery": discovery,
            "selected_products": [{"name": "CEO Daily", "category": "Newsletters"}],
        },
        gtm=gtm,
        inventory=inventory,
    )
    assert confirm["status"] == "ok", confirm
    mix = sum(p["price"] for p in confirm["confirmed_products"])
    deck_schema = _align_budgets(confirm["deck_schema"], mix)
    return deck_schema, confirm


def _verify_deck(prs: Presentation) -> list[str]:
    blob = _slide_text(prs)
    errors: list[str] = []
    for token in LEFTOVER_TOKENS:
        if token in blob:
            errors.append(f"leftover token {token!r}")
    if "CEO Daily" not in blob and "CEO DAILY" not in blob.upper():
        errors.append("CEO Daily clone missing")
    if not re.search(r"\$5,000", blob):
        errors.append("investment budget not filled")
    return errors


def _run_mock_ai(deck_schema: dict) -> int:
    from ingestion.placeholder_fills import apply_placeholders

    gen = DeckGenerator()
    schema = DeckSchema.model_validate(deck_schema)
    prs = gen.assemble_skeleton(schema, template_url=None)
    audience = gen._get_audience_data()
    warnings = apply_placeholders(
        prs,
        schema,
        audience=audience,
        logo_bytes=MINIMAL_PNG,
        ai=mock_placeholder_ai(),
    )
    errors = _verify_deck(prs)
    if warnings:
        print("warnings:", warnings)
    if errors:
        print("SMOKE FAILED:")
        for err in errors:
            print(" -", err)
        return 1
    print(
        "SMOKE OK (ideation→confirm→build, mock AI):",
        len(prs.slides),
        "slides",
    )
    return 0


def _run_live_build(deck_schema: dict) -> int:
    gen = DeckGenerator()
    schema = DeckSchema.model_validate(deck_schema)
    result = gen.build(schema, logo_bytes=MINIMAL_PNG)
    print("build result:", {k: result[k] for k in result if k != "download_url"})
    resp = requests.get(result["download_url"], timeout=60)
    resp.raise_for_status()
    prs = Presentation(io.BytesIO(resp.content))
    errors = _verify_deck(prs)
    if errors:
        print("SMOKE FAILED:")
        for err in errors:
            print(" -", err)
        return 1
    print("SMOKE OK (live build):", result["slide_count"], "slides")
    return 0


def main() -> int:
    load_dotenv()
    if not os.environ.get("S3_SNAPSHOT_BUCKET"):
        print("SKIP: S3_SNAPSHOT_BUCKET not set (copy .env.example → .env)")
        return 0

    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--mock-ai",
        action="store_true",
        help="Assemble + placeholder fills with mock Claude",
    )
    args = parser.parse_args()

    deck_schema, confirm = _run_selection_chain()
    print("confirm:", confirm["status"], confirm.get("confirmed_products"))
    if args.mock_ai:
        return _run_mock_ai(deck_schema)
    return _run_live_build(deck_schema)


if __name__ == "__main__":
    sys.exit(main())
