"""C2 placeholder primitives and variant select."""

from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from ingestion.placeholders import (
    delete_unused_variants,
    select_audience_variant,
    select_program_variant,
)
from ingestion.pptx_tools import (
    APOS,
    CLIENT_NAME_POSSESSIVE_TOKEN,
    CLIENT_NAME_TOKEN,
    LOGO_TOKEN,
    insert_logo,
    replace_first_token,
    replace_token,
)
from tests.fortuneai_placeholder_fixture import (
    HISTORY_BODY,
    MINIMAL_PNG,
    WHY_FORTUNE_STOCK,
    build_fortuneai_fixture_prs,
    fortuneai_fixture_bytes,
)


def _slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            texts.append(para.text)
    return texts


def _embedded_images(slide) -> list:
    images = []
    for shape in slide.shapes:
        try:
            images.append(shape.image)
        except AttributeError:
            continue
    return images


def test_fixture_has_19_slides_and_named_tokens():
    prs = build_fortuneai_fixture_prs()
    assert len(prs.slides) == 19
    intro = " ".join(_slide_texts(prs.slides[0]))
    assert "[TITLE]" in intro
    assert "[DATE]" in intro
    history = " ".join(_slide_texts(prs.slides[2]))
    assert CLIENT_NAME_TOKEN in history
    assert WHY_FORTUNE_STOCK in " ".join(_slide_texts(prs.slides[1]))
    thanks = " ".join(_slide_texts(prs.slides[18]))
    assert "[DATE]" in thanks
    assert "Thank you!" in thanks


def test_replace_token_swaps_history_client_name():
    prs = build_fortuneai_fixture_prs()
    history = prs.slides[2]
    hits = replace_token(history, CLIENT_NAME_TOKEN, "Acme Corp")
    blob = " ".join(_slide_texts(history))
    assert hits >= 1
    assert "Acme Corp" in blob
    assert CLIENT_NAME_TOKEN not in blob
    assert f"Acme Corp{APOS}s opportunity" in blob


def test_replace_token_leaves_why_fortune_unchanged():
    prs = build_fortuneai_fixture_prs()
    why = prs.slides[1]
    before = _slide_texts(why)
    hits = replace_token(why, CLIENT_NAME_TOKEN, "Acme Corp")
    assert hits == 0
    assert _slide_texts(why) == before
    assert WHY_FORTUNE_STOCK in " ".join(before)


def test_replace_token_leaves_unmatched_text_alone():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.paragraphs[0].add_run().text = "Fortune reaches 42 million people"
    hits = replace_token(slide, CLIENT_NAME_TOKEN, "Acme Corp")
    assert hits == 0
    assert "Fortune reaches 42 million people" in " ".join(_slide_texts(slide))


def test_replace_token_audience_possessive_curly_apostrophe():
    prs = build_fortuneai_fixture_prs()
    audience = prs.slides[4]  # 2-card page
    hits = replace_token(audience, CLIENT_NAME_POSSESSIVE_TOKEN, f"Acme Corp{APOS}s")
    blob = " ".join(_slide_texts(audience))
    assert hits >= 1
    assert CLIENT_NAME_POSSESSIVE_TOKEN not in blob
    assert f"OVERDELIVERS Acme Corp{APOS}s TARGET" in blob


def test_replace_token_empty_token_raises():
    prs = build_fortuneai_fixture_prs()
    with pytest.raises(ValueError, match="non-empty"):
        replace_token(prs.slides[0], "", "x")


def test_replace_token_literal_product_type():
    prs = build_fortuneai_fixture_prs()
    program = prs.slides[9]  # 2-box program
    hits = replace_token(program, "PRODUCT TYPE", "Editorial Alignment")
    blob = " ".join(_slide_texts(program))
    assert hits == 2
    assert "PRODUCT TYPE" not in blob
    assert blob.count("Editorial Alignment") == 2


def test_replace_first_token_only_swaps_one_product_type():
    prs = build_fortuneai_fixture_prs()
    program = prs.slides[9]
    assert replace_first_token(program, "PRODUCT TYPE", "Editorial Alignment") == 1
    blob = " ".join(_slide_texts(program))
    assert blob.count("Editorial Alignment") == 1
    assert blob.count("PRODUCT TYPE") == 1


def test_insert_logo_lands_on_intro_and_thanks():
    prs = build_fortuneai_fixture_prs()
    insert_logo(prs.slides[0], MINIMAL_PNG)
    insert_logo(prs.slides[18], MINIMAL_PNG)
    intro_imgs = _embedded_images(prs.slides[0])
    thanks_imgs = _embedded_images(prs.slides[18])
    assert len(intro_imgs) == 1
    assert len(thanks_imgs) == 1
    assert intro_imgs[0].blob.startswith(b"\x89PNG")
    assert thanks_imgs[0].blob.startswith(b"\x89PNG")
    assert LOGO_TOKEN not in " ".join(_slide_texts(prs.slides[0]))
    assert LOGO_TOKEN not in " ".join(_slide_texts(prs.slides[18]))


def test_insert_logo_missing_placeholder_fails_loud():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank, no picture ph
    with pytest.raises(ValueError, match="logo placeholder not found"):
        insert_logo(slide, MINIMAL_PNG)


def test_insert_logo_garbage_bytes_fail_loud():
    prs = build_fortuneai_fixture_prs()
    with pytest.raises(ValueError, match="unreadable"):
        insert_logo(prs.slides[0], b"not-an-image")
    with pytest.raises(ValueError, match="empty"):
        insert_logo(prs.slides[0], b"")


def test_fortuneai_fixture_bytes_roundtrip():
    data = fortuneai_fixture_bytes()
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 19
    assert CLIENT_NAME_TOKEN in HISTORY_BODY


def _count_on_slide(slide, token: str) -> int:
    return " ".join(_slide_texts(slide)).count(token)


def test_select_audience_3_keeps_3_card_index():
    choice = select_audience_variant(3)
    assert choice.keep_index == 5
    assert choice.size == 3
    assert choice.warning is None


def test_select_audience_lt_2_raises():
    with pytest.raises(ValueError, match="at least 2"):
        select_audience_variant(1)
    with pytest.raises(ValueError, match="at least 2"):
        select_audience_variant(0)


def test_select_audience_gt_6_keeps_6_card_and_warns():
    choice = select_audience_variant(8)
    assert choice.keep_index == 8
    assert choice.size == 6
    assert choice.warning is not None
    assert "top 6" in choice.warning


def test_select_program_1_category_keeps_2_box_page():
    choice = select_program_variant(1)
    assert choice.keep_index == 9
    assert choice.size == 2
    assert choice.warning is None


def test_select_program_gt_4_keeps_4_box_and_warns():
    choice = select_program_variant(5)
    assert choice.keep_index == 11
    assert choice.size == 4
    assert choice.warning is not None


def test_delete_unused_audience_3_segments_drops_other_card_pages():
    prs = build_fortuneai_fixture_prs()
    audience = select_audience_variant(3)
    program = select_program_variant(2)
    delete_unused_variants(
        prs, audience_keep=audience.keep_index, program_keep=program.keep_index
    )
    # 19 - 4 unused audience - 2 unused program = 13
    assert len(prs.slides) == 13
    audience_pages = [
        s for s in prs.slides if "[AUDIENCE TITLE]" in " ".join(_slide_texts(s))
    ]
    assert len(audience_pages) == 1
    assert _count_on_slide(audience_pages[0], "[AUDIENCE SEGMENT]") == 3
    program_pages = [
        s for s in prs.slides if "PROGRAM OVERVIEW" in " ".join(_slide_texts(s))
    ]
    assert len(program_pages) == 1
    assert _count_on_slide(program_pages[0], "PRODUCT TYPE") == 2
    # Spine around the kept variants is intact.
    assert WHY_FORTUNE_STOCK in " ".join(_slide_texts(prs.slides[1]))
    assert CLIENT_NAME_TOKEN in " ".join(_slide_texts(prs.slides[2]))
    titles = [
        s.shapes.title.text
        for s in prs.slides
        if s.shapes.title is not None and s.shapes.title.text
    ]
    assert "High-Impact Media" in titles
    assert "Thank you!" in " ".join(_slide_texts(prs.slides[-1]))


def test_delete_unused_1_category_keeps_2_box_program():
    prs = build_fortuneai_fixture_prs()
    audience = select_audience_variant(2)
    program = select_program_variant(1)
    delete_unused_variants(
        prs, audience_keep=audience.keep_index, program_keep=program.keep_index
    )
    program_pages = [
        s for s in prs.slides if "PROGRAM OVERVIEW" in " ".join(_slide_texts(s))
    ]
    assert len(program_pages) == 1
    assert _count_on_slide(program_pages[0], "PRODUCT TYPE") == 2
    assert _count_on_slide(program_pages[0], "Product description.") == 2
