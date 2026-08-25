"""Clone must produce a package PowerPoint opens without repairing.

Every regression here was a real defect in a deck the sales team opened: a
dangling ``r:id`` from a hyperlink, two media parts fighting over one partname,
duplicate shape ids, and product slides silently restyled by the wrong layout.
"""

import collections
import io
import os
import re
import zipfile
from xml.etree import ElementTree as ET

import pytest
from lxml import etree
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches

from ingestion import pptx_tools
from ingestion.generator import _PML_NS, _RELS_NS, DeckGenerator
from tests.fortuneai_placeholder_fixture import MINIMAL_PNG

_SLIDE_RE = re.compile(r"ppt/(slides|slideLayouts|slideMasters)/[a-zA-Z]+\d+\.xml$")


def _save(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _dangling_rids(data: bytes) -> list[tuple[str, list[str]]]:
    """Parts referencing an r:id their .rels file does not define."""
    zf = zipfile.ZipFile(io.BytesIO(data))
    names = zf.namelist()
    bad = []
    for name in filter(_SLIDE_RE.match, names):
        rels_name = f"{os.path.dirname(name)}/_rels/{os.path.basename(name)}.rels"
        declared = (
            {r.get("Id") for r in ET.fromstring(zf.read(rels_name))}
            if rels_name in names
            else set()
        )
        used = set(re.findall(r'r:(?:id|embed|link)="([^"]+)"', zf.read(name).decode()))
        if used - declared:
            bad.append((name, sorted(used - declared)))
    return bad


def _missing_targets(data: bytes) -> list[tuple[str, str]]:
    zf = zipfile.ZipFile(io.BytesIO(data))
    names = zf.namelist()
    missing = []
    for name in (n for n in names if n.endswith(".rels")):
        base = os.path.dirname(os.path.dirname(name))
        for rel in ET.fromstring(zf.read(name)):
            if rel.get("TargetMode") == "External":
                continue
            target = os.path.normpath(os.path.join(base, rel.get("Target")))
            if target not in names:
                missing.append((name, rel.get("Target")))
    return missing


def _duplicate_partnames(data: bytes) -> dict[str, int]:
    names = zipfile.ZipFile(io.BytesIO(data)).namelist()
    return {n: c for n, c in collections.Counter(names).items() if c > 1}


def _shape_ids(slide) -> list[int]:
    return [int(el.get("id")) for el in slide._element.iter(qn("p:cNvPr"))]


def _png(color: bytes = b"\x00") -> io.BytesIO:
    return io.BytesIO(MINIMAL_PNG)


def _source_with_pictures(count: int = 4, *, hyperlink: bool = True):
    """Source deck whose slide carries several pictures and an external link."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for i in range(count):
        pic = slide.shapes.add_picture(_png(), Inches(i), Inches(1), Inches(1), Inches(1))
        if hyperlink and i == 0:
            pic.click_action.hyperlink.address = "https://example.fortune.com/x"
    return prs


def test_clone_carries_non_image_relationships():
    """Hyperlink rels used to be dropped, leaving an r:id PowerPoint calls corrupt."""
    source = _source_with_pictures()
    target = Presentation()

    DeckGenerator._clone_slide(source, 0, target)

    assert _dangling_rids(_save(target)) == []


def test_remap_rids_handles_ids_that_swap_places():
    """The old sequential string replace collapsed rIds whose ranges overlapped.

    Mapping rId2->rId3 first rewrote the tag that already said rId3, so the
    second mapping caught both and two pictures ended up sharing one image.
    """
    xml = (
        f'<p:pics xmlns:p="{_PML_NS}" xmlns:r="{_RELS_NS}">'
        '<p:a r:embed="rId2"/><p:b r:embed="rId3"/><p:c r:id="rId4"/>'
        "</p:pics>"
    )
    element = parse_xml(xml)

    DeckGenerator._remap_rids(element, {"rId2": "rId3", "rId3": "rId2", "rId4": "rId9"})

    rewritten = etree.tostring(element).decode()
    assert re.findall(r'r:(?:embed|id)="([^"]+)"', rewritten) == ["rId3", "rId2", "rId9"]


def test_cloned_pictures_only_reference_declared_relationships():
    source = _source_with_pictures(hyperlink=False)
    target = Presentation()

    clone = DeckGenerator._clone_slide(source, 0, target)

    embeds = set(re.findall(r'r:embed="([^"]+)"', clone._element.xml))
    assert embeds
    assert embeds <= set(clone.part.rels)


def test_clone_does_not_duplicate_media_partnames():
    """Product decks and FortuneAI both ship ppt/media/image1.png."""
    source = _source_with_pictures(count=1, hyperlink=False)
    target = Presentation()
    target.slides.add_slide(target.slide_layouts[5]).shapes.add_picture(
        _png(), Inches(0), Inches(0), Inches(1), Inches(1)
    )

    DeckGenerator._clone_slide(source, 0, target)
    data = _save(target)

    assert _duplicate_partnames(data) == {}
    assert _missing_targets(data) == []


def test_dedup_shape_ids_never_creates_a_duplicate():
    """Replacement ids used to collide with shapes further down the same slide."""
    target = Presentation()
    slide = target.slides.add_slide(target.slide_layouts[5])
    for i in range(4):
        slide.shapes.add_textbox(Inches(i), Inches(1), Inches(1), Inches(1))
    # Force the shape the old algorithm would hand out (max+1) to already exist
    # later in the tree, which silently produced two shapes with one id.
    elements = list(slide._element.iter(qn("p:cNvPr")))
    for el, sid in zip(elements, ["1", "1", "1", "3", "2"]):
        el.set("id", sid)

    DeckGenerator._dedup_shape_ids(slide, target)

    ids = _shape_ids(slide)
    assert len(ids) == len(set(ids)), f"duplicate cNvPr ids: {ids}"


def test_clone_imports_source_layout_master_and_theme():
    """Placeholders inherit from the layout, so the clone must keep its own."""
    source = Presentation()
    source_slide = source.slides.add_slide(source.slide_layouts[1])
    source_slide.shapes.title.text = "CEO DAILY"
    target = Presentation()
    before = len(target.slide_masters)

    clone = DeckGenerator._clone_slide(source, 0, target)

    assert clone.slide_layout.name == source_slide.slide_layout.name
    assert len(target.slide_masters) == before + 1
    assert _dangling_rids(_save(target)) == []


def test_repeated_clones_share_one_imported_master():
    """Cloning five products off one deck must not bolt on five masters."""
    source = Presentation()
    for _ in range(3):
        source.slides.add_slide(source.slide_layouts[1])
    target = Presentation()
    before = len(target.slide_masters)

    for idx in range(3):
        DeckGenerator._clone_slide(source, idx, target)

    assert len(target.slide_masters) == before + 1


def test_imported_master_is_pruned_to_the_used_layout():
    """Importing an untrimmed master would drag a 300-layout product deck along."""
    source = Presentation()
    source.slides.add_slide(source.slide_layouts[1])
    target = Presentation()

    clone = DeckGenerator._clone_slide(source, 0, target)

    imported_master = clone.slide_layout.slide_master
    assert len(imported_master.slide_layouts) == 1


_P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"


def _sectioned_deck(section_map: dict[str, int]):
    """Deck with ``section_map`` of section name -> slide count, in slide order."""
    prs = Presentation()
    sections = []
    for name, count in section_map.items():
        ids = []
        for _ in range(count):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            ids.append(prs.slides._sldIdLst[-1].get("id"))
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
        entries = "".join(f'<p14:sldId id="{i}"/>' for i in ids)
        sections.append(
            f'<p14:section name="{name}" id="{{{name}}}">'
            f"<p14:sldIdLst>{entries}</p14:sldIdLst></p14:section>"
        )
    prs.part._element.append(
        parse_xml(
            f'<p:extLst xmlns:p="{_PML_NS}">'
            '<p:ext uri="{521415D9-36F7-43E2-AB2F-B90AF26B5E84}">'
            f'<p14:sectionLst xmlns:p14="{_P14_NS}">{"".join(sections)}</p14:sectionLst>'
            "</p:ext></p:extLst>"
        )
    )
    return prs


def _sections(prs) -> dict[str, list[str]]:
    lst = pptx_tools._section_lst(prs)
    if lst is None:
        return {}
    return {
        s.get("name"): [e.get("id") for e in s.iter(f"{{{_P14_NS}}}sldId")] for s in lst
    }


def test_delete_slide_drops_its_own_section_entry():
    """Deletions run after assembly too (unused audience pages), so delete must self-clean."""
    prs = _sectioned_deck({"KEEP": 1, "AUDIENCE": 3})
    doomed = prs.slides._sldIdLst[2].get("id")

    pptx_tools.delete_slide(prs, 2)

    assert doomed not in _sections(prs)["AUDIENCE"]
    live = {sld.get("id") for sld in prs.slides._sldIdLst}
    assert {s for ids in _sections(prs).values() for s in ids} <= live


def test_sync_sections_drops_references_to_deleted_slides():
    """A section pointing at a removed slide is what PowerPoint 'repairs'."""
    prs = _sectioned_deck({"KEEP": 2, "DIVIDERS": 2})
    pptx_tools.delete_slide(prs, 2)

    pptx_tools.sync_sections(prs)

    live = {sld.get("id") for sld in prs.slides._sldIdLst}
    referenced = [sid for ids in _sections(prs).values() for sid in ids]
    assert set(referenced) <= live
    assert sorted(referenced) == sorted(live)


def test_sync_sections_removes_a_section_left_with_no_slides():
    """Deleting the stock Print page used to leave an empty 'Print Page' section."""
    prs = _sectioned_deck({"KEEP": 1, "PRINT": 1})
    pptx_tools.delete_slide(prs, 1)

    pptx_tools.sync_sections(prs)

    assert "PRINT" not in _sections(prs)


def test_sync_sections_files_new_slides_under_the_preceding_section():
    prs = _sectioned_deck({"INTRO": 1, "OFFERINGS": 1, "THANKS": 1})
    source = Presentation()
    source.slides.add_slide(source.slide_layouts[5])
    DeckGenerator._clone_slide(source, 0, prs)
    DeckGenerator._insert_slide_at(prs, 2)
    clone_id = prs.slides._sldIdLst[2].get("id")

    pptx_tools.sync_sections(prs)

    assert clone_id in _sections(prs)["OFFERINGS"]


def test_sync_sections_is_a_no_op_without_sections():
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])

    pptx_tools.sync_sections(prs)

    assert _sections(prs) == {}


@pytest.mark.parametrize("slide_idx", [0, 1])
def test_cloned_package_has_no_missing_relationship_targets(slide_idx):
    source = _source_with_pictures()
    source.slides.add_slide(source.slide_layouts[1]).shapes.title.text = "Second"
    target = Presentation()

    DeckGenerator._clone_slide(source, slide_idx, target)

    assert _missing_targets(_save(target)) == []
