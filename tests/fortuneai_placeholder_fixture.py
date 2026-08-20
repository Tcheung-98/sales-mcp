"""CI FortuneAI-shaped fixture for placeholder tests (no live SharePoint/S3).

Slide roles match the Workflow spine: intro, narrative variants, category
dividers, investment, thanks. Intro/thanks use a picture placeholder labeled
[LOGO]; other slides use text boxes with named tokens.
"""

from __future__ import annotations

import io
from unittest.mock import MagicMock

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from ingestion.audience_data import AudienceData, AudienceRow
from ingestion.pptx_tools import APOS, CLIENT_NAME_POSSESSIVE_TOKEN, CLIENT_NAME_TOKEN, LOGO_TOKEN

# 1x1 PNG — valid image for insert_logo tests (no SharePoint).
MINIMAL_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
    "0000000a49444154789c63000100000500010d0a2db40000000049454e44ae426082"
)


def sample_audience_data() -> AudienceData:
    """Small Audience Data index for CI (no live GTM xlsx)."""
    return AudienceData(
        [
            AudienceRow("Chief Executive Officer", "1.1M", "154"),
            AudienceRow("C-suite", "3.6M", "172"),
            AudienceRow("Chief Financial Officer", "422K", "146"),
            AudienceRow("Chief Information Officer", "52K", "207"),
            AudienceRow("Chief Technology Officer", "970K", "255"),
            AudienceRow("Chief Data Officer", "299K", "219"),
            AudienceRow("Active Investor", "2.3M", "118"),
            AudienceRow("Wealthy (HNW)", "697K", "161"),
        ]
    )

# Valid Opportunity body for mocked AI fills in placeholder tests (~85 words).
SAMPLE_OPPORTUNITY_BODY = (
    "Enterprise buyers are skeptical of vendor claims and generic advertising noise. "
    "Acme Corp can stand out by aligning with editorial environments CFOs trust "
    "for business insight daily. Finance leaders seek credible partners who "
    "understand modernization without hype in complex markets today. Engaging them "
    "requires authority, context, and relevance rather than interruptive formats "
    "alone across channels. Fortune delivers authority through newsletters, live "
    "experiences, and premium brand storytelling trusted by executives worldwide. "
    "Acme can lead the category conversation by showing up where decisions happen. "
    "Fortune connects your message to the executives shaping the future of business."
)

SAMPLE_PROGRAM_BLURB = (
    "Fortune newsletters connect Acme with decision-makers through trusted "
    "weekly editorial environments and strategic insight daily."
)


WHY_FORTUNE_STOCK = "FORTUNE POWERS THE LEADING MINDS IN BUSINESS"
HISTORY_BODY = (
    f"In a fractured media landscape, that trust is Fortune{APOS}s edge — "
    f"and {CLIENT_NAME_TOKEN}{APOS}s opportunity to show up on a platform "
    "trusted by the leaders who are changing the world."
)


def mock_placeholder_ai() -> MagicMock:
    ai = MagicMock()
    ai.intro_title.return_value = "ACME CORP ENTERPRISE PARTNERSHIP"
    ai.opportunity_header.return_value = "Lead With Confidence Today"
    ai.opportunity_body.return_value = SAMPLE_OPPORTUNITY_BODY
    ai.audience_title.return_value = "Reach enterprise leaders"
    ai.program_blurb.return_value = SAMPLE_PROGRAM_BLURB
    return ai


def _add_textbox(slide, text: str, *, top_in: float = 1.0) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top_in), Inches(8), Inches(1.2))
    box.text_frame.paragraphs[0].add_run().text = text


def _picture_placeholder(slide):
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                return ph
        except (ValueError, AttributeError):
            continue
    raise RuntimeError("fixture layout has no picture placeholder")


def _intro_or_thanks(prs: Presentation, *, thanks: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption
    title = slide.placeholders[0]
    title.text = "Thank you!" if thanks else "[TITLE]"
    pic = _picture_placeholder(slide)
    if pic.has_text_frame:
        para = pic.text_frame.paragraphs[0]
        if para.runs:
            para.runs[0].text = LOGO_TOKEN
        else:
            para.text = LOGO_TOKEN
    body = slide.placeholders[2]
    body.text = "[DATE]"


def build_fortuneai_fixture_prs() -> Presentation:
    """19-slide spine with named tokens on the Workflow roles."""
    prs = Presentation()
    _intro_or_thanks(prs, thanks=False)

    why = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    why.shapes.title.text = WHY_FORTUNE_STOCK
    _add_textbox(why, "WHY fortune")

    history = prs.slides.add_slide(prs.slide_layouts[5])
    history.shapes.title.text = "TRUST IS THE ULTIMATE COMPETITIVE ADVANTAGE"
    _add_textbox(history, HISTORY_BODY)

    opp = prs.slides.add_slide(prs.slide_layouts[5])
    opp.shapes.title.text = "[HEADER]"
    _add_textbox(opp, "[BODY]")

    for cards in (2, 3, 4, 5, 6):
        aud = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        _add_textbox(
            aud,
            f"FORTUNE OVERDELIVERS {CLIENT_NAME_POSSESSIVE_TOKEN} TARGET AUDIENCE",
            top_in=0.3,
        )
        _add_textbox(aud, "[AUDIENCE TITLE]", top_in=0.8)
        for i in range(cards):
            _add_textbox(
                aud,
                "[AUDIENCE SEGMENT]|[REACH]|[INDEX]",
                top_in=1.3 + i * 0.4,
            )

    for boxes in (2, 3, 4):
        prog = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(prog, "PROGRAM OVERVIEW", top_in=0.2)
        for i in range(boxes):
            _add_textbox(prog, "PRODUCT TYPE", top_in=0.8 + i * 0.7)
            _add_textbox(prog, "Product description.", top_in=1.1 + i * 0.7)

    for name in (
        "High-Impact Media",
        "Editorial Alignment",
        "Premium Video",
        "Print",
        "Branded Content",
    ):
        div = prs.slides.add_slide(prs.slide_layouts[5])
        div.shapes.title.text = name

    inv = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(inv, "[BUDGET]", top_in=0.4)
    cat = inv.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(8), Inches(1.6))
    cat.text_frame.paragraphs[0].add_run().text = "[PRODUCT CATEGORY]"
    price_para = cat.text_frame.add_paragraph()
    price_para.add_run().text = "[PRODUCT + PRODUCT PRICE] "

    _intro_or_thanks(prs, thanks=True)
    return prs


def fortuneai_fixture_bytes() -> bytes:
    prs = build_fortuneai_fixture_prs()
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
