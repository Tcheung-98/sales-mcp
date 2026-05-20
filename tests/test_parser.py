import io

from pptx import Presentation

from ingestion.models import Tags
from ingestion.parser import content_hash, parse_pptx


def make_pptx(slides: list[dict]) -> bytes:
    """Build an in-memory .pptx from a list of slide specs. No files, no network.
    Each spec may include: title, body, notes, layout (int, default 1).
    Use layout=6 for a blank slide with no title placeholder.
    """
    prs = Presentation()
    for spec in slides:
        layout = prs.slide_layouts[spec.get("layout", 1)]
        slide = prs.slides.add_slide(layout)
        if spec.get("title") and slide.shapes.title:
            slide.shapes.title.text = spec["title"]
        if spec.get("body") and len(slide.placeholders) > 1:
            slide.placeholders[1].text = spec["body"]
        if spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = spec["notes"]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


SAMPLE_TAGS = Tags(industry="Tech")

# happy path
def test_single_slide_returns_one_sliderow():
    pptx_bytes = make_pptx([
        {"title": "Overview", "body": "Q1 Revenue up 15%", "notes": "Discuss YoY growth"},
    ])

    rows = parse_pptx(
        pptx_bytes=pptx_bytes,
        deck_id="deck-001",
        source_path="/TechMedia/acme.pptx",
        ingested_at="2026-05-18T10:00:00",
        tags=SAMPLE_TAGS,
    )

    assert len(rows) == 1
    row = rows[0]
    assert row.deck_id == "deck-001"
    assert row.slide_number == 1
    assert row.title == "Overview"
    assert "Q1 Revenue up 15%" in row.body_text
    assert row.speaker_notes == "Discuss YoY growth"
    assert row.tags.industry == "Tech"
    assert isinstance(row.body_text, list)
    print('test_single_slide_returns_one_sliderow passed!')

def _parse(pptx_bytes):
    return parse_pptx(
        pptx_bytes=pptx_bytes,
        deck_id="deck-001",
        source_path="/TechMedia/acme.pptx",
        ingested_at="2026-05-18T10:00:00",
        tags=SAMPLE_TAGS,
    )


# Zero-slide deck returns empty list
def test_zero_slide_returns_empty_list():
    rows = _parse(make_pptx([]))
    assert rows == []


# Slide with no title placeholder returns title=None
def test_no_title_returns_none():
    rows = _parse(make_pptx([{"layout": 6}]))  # blank layout has no title placeholder
    assert len(rows) == 1
    assert rows[0].title is None


# Slide with no speaker notes returns empty string, never None
def test_no_speaker_notes_returns_empty_string():
    rows = _parse(make_pptx([{"title": "Slide", "body": "Content"}]))
    assert rows[0].speaker_notes == ""
    assert rows[0].speaker_notes is not None


# Non-breaking spaces and vertical tabs are cleaned from body text
def test_special_chars_are_cleaned():
    pptx_bytes = make_pptx([{"title": "Title", "body": "Revenue\xa0up\x0b15%"}])
    rows = _parse(pptx_bytes)
    assert rows[0].body_text != []
    assert all("\xa0" not in line and "\x0b" not in line for line in rows[0].body_text)


# Multiple slides are numbered sequentially starting at 1
def test_multiple_slides_correct_slide_numbers():
    rows = _parse(make_pptx([
        {"title": "Slide One"},
        {"title": "Slide Two"},
        {"title": "Slide Three"},
    ]))
    assert len(rows) == 3
    assert [r.slide_number for r in rows] == [1, 2, 3]


# content_hash is deterministic — same bytes always produce the same hash
def test_content_hash_is_deterministic():
    pptx_bytes = make_pptx([{"title": "Test"}])
    assert content_hash(pptx_bytes) == content_hash(pptx_bytes)


# content_hash differs for different file contents
def test_content_hash_differs_for_different_content():
    a = make_pptx([{"title": "Deck A"}])
    b = make_pptx([{"title": "Deck B"}])
    assert content_hash(a) != content_hash(b)