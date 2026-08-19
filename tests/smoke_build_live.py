"""Live-ish smoke test: real S3 FortuneAI + GTM, mocked Claude, local PPTX check.

  PYTHONPATH=. uv run python tests/smoke_build_live.py --mock-ai

Use without --mock-ai only when Anthropic and logo HTTPS are reachable.
"""

from __future__ import annotations

import argparse
import io
import re
import sys
from dotenv import load_dotenv
from pptx import Presentation

from ingestion.generator import DeckGenerator
from ingestion.schema import DeckSchema, Product
from tests.fortuneai_placeholder_fixture import MINIMAL_PNG, mock_placeholder_ai

LEFTOVER_TOKENS = (
    "[TITLE]",
    "[HEADER]",
    "[BODY]",
    "[AUDIENCE TITLE]",
    "Product description.",
    "[DATE]",
    "[LOGO]",
    "[client name]",
)


def _slide_text(prs: Presentation) -> str:
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                parts.append(para.text or "")
    return "\n".join(parts)


def _schema() -> DeckSchema:
    return DeckSchema(
        company_name="Acme Corp",
        industry="Technology",
        budgets=[{"amount": 50_000}],
        flight_dates={"start": "2026-09-01", "end": "2026-12-31"},
        campaign_goal="Drive consideration among enterprise buyers",
        targeting_details=(
            "Chief Executive Officer, C-suite, Chief Financial Officer"
        ),
        kpis=["Awareness", "Engagement"],
        kpi_details="Lift brand awareness 10%; engagement rate above benchmark",
        campaign_narrative="Acme helps mid-market CFOs modernize finance ops",
        preferred_platforms_products=["Newsletters", "Branded Content"],
        additional_rfp_details="Prefer Q4 flight",
        client_logo="https://example.com/logo.png",
        confirmed_products=[
            Product(
                name="CEO Daily",
                cadence="weekly",
                price=50_000,
                category="Newsletter",
            )
        ],
    )


def _verify_prs(prs: Presentation, *, slide_count: int) -> list[str]:
    blob = _slide_text(prs)
    titles = [
        s.shapes.title.text
        for s in prs.slides
        if s.shapes.title is not None and s.shapes.title.text
    ]
    errors: list[str] = []
    if len(prs.slides) != slide_count:
        errors.append(f"expected {slide_count} slides, got {len(prs.slides)}")
    for token in LEFTOVER_TOKENS:
        if token in blob:
            errors.append(f"leftover token {token!r}")
    if "CEO DAILY" not in titles and "CEO Daily" not in blob:
        errors.append("CEO Daily clone missing")
    normalized = re.sub(r"[\x0b\r\n]+", " ", blob)
    if "FORTUNE POWERS THE" not in normalized or "LEADING MINDS IN BUSINESS" not in normalized:
        errors.append("Why Fortune stock copy missing")
    if not re.search(r"\$50,000", blob):
        errors.append("investment budget not filled")
    if "[AUDIENCE SEGMENT]" in blob:
        errors.append("unfilled audience segment tokens")
    return errors


def _run_mock_ai() -> int:
    from ingestion.placeholder_fills import apply_placeholders

    gen = DeckGenerator()
    schema = _schema()
    prs = gen.assemble_skeleton(schema, template_url=None)
    audience = gen._get_audience_data()
    warnings = apply_placeholders(
        prs,
        schema,
        audience=audience,
        logo_bytes=MINIMAL_PNG,
        ai=mock_placeholder_ai(),
    )
    errors = _verify_prs(prs, slide_count=10)
    if warnings:
        print("warnings:", warnings)
    if errors:
        print("SMOKE FAILED:")
        for err in errors:
            print(" -", err)
        return 1
    print("SMOKE OK (mock AI):", len(prs.slides), "slides on real FortuneAI template")
    return 0


def _run_live() -> int:
    gen = DeckGenerator()
    schema = _schema()
    result = gen.build(schema, logo_bytes=MINIMAL_PNG)
    print("build result:", {k: result[k] for k in result if k != "download_url"})
    resp = requests.get(result["download_url"], timeout=60)
    resp.raise_for_status()
    prs = Presentation(io.BytesIO(resp.content))
    errors = _verify_prs(prs, slide_count=result["slide_count"])
    if result.get("warnings"):
        print("warnings:", result["warnings"])
    if errors:
        print("SMOKE FAILED:")
        for err in errors:
            print(" -", err)
        return 1
    print("SMOKE OK (live Claude):", result["slide_count"], "slides")
    return 0


def main() -> int:
    load_dotenv()
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--mock-ai",
        action="store_true",
        help="Use real S3 template/GTM but mock Claude (no outbound API)",
    )
    args = parser.parse_args()
    if args.mock_ai:
        return _run_mock_ai()
    return _run_live()


if __name__ == "__main__":
    sys.exit(main())
