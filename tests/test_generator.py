import io
from unittest.mock import MagicMock, patch

import anthropic
import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from ingestion.generator import DeckGenerator
from ingestion.gtm_product_map import GtmProductMap, ProductSlideRef
from ingestion.pptx_tools import apply_replacements
from ingestion.schema import DeckSchema, Product
from tests.fortuneai_placeholder_fixture import (
    MINIMAL_PNG,
    fortuneai_fixture_bytes,
    sample_audience_data,
)


def _blank_bytes(slide_count: int = 1) -> bytes:
    """Minimal blank template with enough slides for _clone_slide(source, 0, ...)."""
    prs = Presentation()
    for _ in range(slide_count):
        prs.slides.add_slide(prs.slide_layouts[1])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _fortuneai_template_bytes(*, extra_tail: int = 0) -> bytes:
    """Minimal FortuneAI spine: 12 narrative + 5 dividers + investment + thank you.

    Divider slides are titled DIV-0 .. DIV-4 so tests can assert which survived.
    """
    prs = Presentation()
    # Slides 1–12: intro + narrative
    for i in range(12):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"SPINE-{i + 1}"
    # Slides 13–17: category dividers
    for i in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"DIV-{i}"
    # Slide 18 investment, 19 thank you
    inv = prs.slides.add_slide(prs.slide_layouts[1])
    inv.shapes.title.text = "INVESTMENT"
    ty = prs.slides.add_slide(prs.slide_layouts[1])
    ty.shapes.title.text = "THANK YOU"
    for i in range(extra_tail):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"EXTRA-{i + 1}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _presentation_with_n_slides(n: int) -> Presentation:
    prs = Presentation()
    for _ in range(n):
        prs.slides.add_slide(prs.slide_layouts[1])
    return prs


def _build_generator() -> DeckGenerator:
    with patch("boto3.client", return_value=MagicMock()):
        generator = DeckGenerator(bucket="test-bucket")
    generator._s3 = MagicMock()
    generator._blank_bytes = _blank_bytes(slide_count=6)
    return generator


def _build_schema(**overrides) -> DeckSchema:
    defaults = dict(
        company_name="Acme Corp",
        industry="Technology",
        budgets=[{"amount": 50_000}],
        flight_dates={"start": "2026-09-01", "end": "2026-12-31"},
        campaign_goal="Drive consideration among enterprise buyers",
        targeting_details=(
            "US enterprise tech decision-makers, Chief Executive Officer, C-suite"
        ),
        kpis=["Awareness", "Engagement"],
        kpi_details="Lift brand awareness 10%; engagement rate above benchmark",
        campaign_narrative="Acme helps mid-market CFOs modernize finance ops",
        preferred_platforms_products=["Newsletters", "Branded Content"],
        additional_rfp_details="Prefer Q4 flight; avoid holiday blackout weeks",
        client_logo="https://example.com/acme-logo.png",
        confirmed_products=[
            Product(
                name="Fortune 500 List",
                cadence="annual",
                price=50_000,
                category="Newsletter",
            )
        ],
    )
    defaults.update(overrides)
    return DeckSchema(**defaults)


_FORTUNEAI_URL = (
    "https://fortune.sharepoint.com/sites/x/FortuneAI_DeckTemplate.pptx"
)


def _slide_with_textbox(text: str):
    """Minimal slide with a text box containing the given text as a single run."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.paragraphs[0].add_run().text = text
    return slide


def test_apply_replacements_swaps_client_name():
    slide = _slide_with_textbox("your brand reaches millions")
    apply_replacements(slide, {"client_name": "Acme Corp"})
    texts = [
        run.text
        for sh in slide.shapes if sh.has_text_frame
        for para in sh.text_frame.paragraphs
        for run in para.runs
    ]
    assert any("Acme Corp" in t for t in texts)
    assert not any("your brand" in t.lower() for t in texts)


def test_apply_replacements_case_insensitive():
    slide = _slide_with_textbox("YOUR COMPANY is the leader")
    apply_replacements(slide, {"client_name": "Acme Corp"})
    texts = [
        run.text
        for sh in slide.shapes if sh.has_text_frame
        for para in sh.text_frame.paragraphs
        for run in para.runs
    ]
    assert any("Acme Corp" in t for t in texts)


def test_apply_replacements_title_overwrites_placeholder():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title, Content
    slide.shapes.title.text = "Old Title"
    apply_replacements(slide, {"title": "New Title"})
    assert slide.shapes.title.text == "New Title"


def test_apply_replacements_leaves_unmatched_text_alone():
    slide = _slide_with_textbox("Fortune reaches 42 million people")
    apply_replacements(slide, {"client_name": "Acme Corp"})
    texts = [
        run.text
        for sh in slide.shapes if sh.has_text_frame
        for para in sh.text_frame.paragraphs
        for run in para.runs
    ]
    assert any("Fortune reaches 42 million people" in t for t in texts)


def test_apply_replacements_noop_on_empty_dict():
    slide = _slide_with_textbox("your brand")
    apply_replacements(slide, {})
    texts = [
        run.text
        for sh in slide.shapes if sh.has_text_frame
        for para in sh.text_frame.paragraphs
        for run in para.runs
    ]
    assert any("your brand" in t for t in texts)


def _rulebook_bytes(text: str = "Rule 1: always outline first.\nRule 2: $750K escalates.") -> bytes:
    doc = Document()
    for line in text.splitlines():
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_load_rulebook_extracts_text():
    generator = _build_generator()
    generator._s3.get_object.return_value = {"Body": io.BytesIO(_rulebook_bytes())}

    text = generator._load_rulebook()

    assert "always outline first" in text
    assert "$750K escalates" in text
    generator._s3.get_object.assert_called_once()


def test_load_rulebook_caches():
    generator = _build_generator()
    generator._s3.get_object.return_value = {"Body": io.BytesIO(_rulebook_bytes())}

    generator._load_rulebook()
    generator._load_rulebook()

    generator._s3.get_object.assert_called_once()


def test_load_pptx_caches_by_key():
    generator = _build_generator()
    generator._s3.get_object.return_value = {"Body": io.BytesIO(_blank_bytes())}

    prs_a1 = generator._load_pptx("corpus/deck_a.pptx")
    prs_a2 = generator._load_pptx("corpus/deck_a.pptx")

    assert prs_a1 is prs_a2
    generator._s3.get_object.assert_called_once()


def _clone_slide_data(source_path: str, slide_number: int = 1, **replacements) -> dict:
    return {
        "action": "clone",
        "source_path": source_path,
        "slide_number": slide_number,
        "replacements": replacements,
    }


def test_build_pptx_correct_slide_count():
    generator = _build_generator()
    source_prs = Presentation(io.BytesIO(_blank_bytes()))
    with patch.object(generator, "_load_pptx", return_value=source_prs):
        slides = [
            {"action": "cover", "title": "Client Growth Plan", "subtitle": "Q3 2026"},
            _clone_slide_data("corpus/deck.pptx"),
        ]
        pptx_bytes = generator._build_pptx(slides)
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 2


def test_build_pptx_title_slide_sets_headline():
    generator = _build_generator()
    slides = [{"action": "cover", "title": "Client Growth Plan", "subtitle": "Q3 2026"}]
    pptx_bytes = generator._build_pptx(slides)
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[0]
    all_text = {
        ph.placeholder_format.idx: ph.text_frame.text
        for ph in slide.placeholders
        if ph.has_text_frame
    }
    assert "Client Growth Plan" in all_text.values()
    assert "Q3 2026" in all_text.values()


def test_build_pptx_clone_action_calls_load_pptx():
    generator = _build_generator()
    source_prs = Presentation(io.BytesIO(_blank_bytes()))
    with patch.object(generator, "_load_pptx", return_value=source_prs) as mock_load:
        slides = [
            {"action": "cover", "title": "T", "subtitle": "S"},
            _clone_slide_data("corpus/product.pptx", slide_number=1),
        ]
        pptx_bytes = generator._build_pptx(slides)
    mock_load.assert_called_once_with("corpus/product.pptx")
    assert len(Presentation(io.BytesIO(pptx_bytes)).slides) == 2


def test_build_pptx_multiple_clone_actions():
    generator = _build_generator()
    source_prs = Presentation(io.BytesIO(_blank_bytes()))
    with patch.object(generator, "_load_pptx", return_value=source_prs):
        slides = [
            {"action": "cover", "title": "T", "subtitle": "S"},
            _clone_slide_data("corpus/deck.pptx"),
            _clone_slide_data("corpus/deck.pptx"),
            _clone_slide_data("corpus/deck.pptx"),
            _clone_slide_data("corpus/deck.pptx"),
        ]
        pptx_bytes = generator._build_pptx(slides)
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 5


def test_call_claude_context_includes_full_coordinates():
    generator = _build_generator()
    generator._rulebook_text = "Rule: always lead with value."
    generator._api_key = "test-key"

    arc = [
        {"slot": 0, "role": "opener", "query": "market opportunity audience"},
        {"slot": 1, "role": "product", "query": "audience reach"},
    ]
    context_by_slot = {
        0: [
            {
                "source_path": "corpus/Fortune_GP_2026.pptx",
                "slide_number": 7,
                "title": "Market Opportunity",
                "body_text": ["$10B TAM", "Fortune reach: 42M"],
            }
        ],
        1: [
            {
                "source_path": "corpus/Fortune_500_2025.pptx",
                "slide_number": 3,
                "title": "Audience",
                "body_text": [],
            }
        ],
    }

    captured: dict = {}

    def fake_create(**kwargs):
        captured["user_msg"] = kwargs["messages"][0]["content"]
        mock_block = MagicMock(spec=anthropic.types.TextBlock)
        mock_block.text = '[{"action": "cover", "title": "T", "subtitle": "S"}]'
        mock_response = MagicMock()
        mock_response.usage.input_tokens = 100
        mock_response.usage.output_tokens = 50
        mock_response.content = [mock_block]
        return mock_response

    with patch("anthropic.Anthropic") as MockAnthropic:
        MockAnthropic.return_value.messages.create.side_effect = fake_create
        generator._call_claude("Test brief", arc, context_by_slot)

    user_msg = captured["user_msg"]
    assert "corpus/Fortune_GP_2026.pptx" in user_msg
    assert "slide: 7" in user_msg
    assert "corpus/Fortune_500_2025.pptx" in user_msg
    assert "slide: 3" in user_msg


def test_delete_slide_removes_slide():
    prs = _presentation_with_n_slides(3)
    DeckGenerator._delete_slide(prs, 1)
    assert len(prs.slides) == 2


def test_delete_slide_correct_slide_removed():
    prs = _presentation_with_n_slides(3)
    titles = ["A", "B", "C"]
    for slide, title in zip(prs.slides, titles):
        slide.shapes.title.text = title
    DeckGenerator._delete_slide(prs, 1)
    remaining = [s.shapes.title.text for s in prs.slides]
    assert remaining == ["A", "C"]


def test_insert_slide_at_moves_to_position():
    prs = _presentation_with_n_slides(3)
    titles = ["A", "B", "C"]
    for slide, title in zip(prs.slides, titles):
        slide.shapes.title.text = title
    source_prs = _presentation_with_n_slides(1)
    source_prs.slides[0].shapes.title.text = "NEW"
    DeckGenerator._clone_slide(source_prs, 0, prs)
    DeckGenerator._insert_slide_at(prs, 1)
    assert len(prs.slides) == 4
    assert prs.slides[1].shapes.title.text == "NEW"
    assert prs.slides[0].shapes.title.text == "A"
    assert prs.slides[2].shapes.title.text == "B"


def test_validate_template_url_allows_sharepoint():
    DeckGenerator._validate_template_url(
        "https://fortune.sharepoint.com/sites/sales/template.pptx"
    )


def test_validate_template_url_rejects_http():
    with pytest.raises(ValueError, match="HTTPS"):
        DeckGenerator._validate_template_url(
            "http://fortune.sharepoint.com/template.pptx"
        )


def test_validate_template_url_rejects_unknown_host():
    with pytest.raises(ValueError, match="host not allowed"):
        DeckGenerator._validate_template_url("https://evil.example.com/template.pptx")


def test_validate_template_url_allows_extra_host(monkeypatch):
    monkeypatch.setenv("TEMPLATE_URL_ALLOWED_HOSTS", "cdn.example.com")
    DeckGenerator._validate_template_url("https://cdn.example.com/template.pptx")


def _product_map_for(*products: Product) -> GtmProductMap:
    """Minimal map so assemble_skeleton can resolve each confirmed product."""
    rows = [
        ProductSlideRef(
            product_name=p.name,
            category={
                "Newsletter": "Newsletters",
                "Digital Media": "Digital Ads/Programmatic",
                "Branded Content": "Branded Content",
                "Print": "Print",
                "Vodcasts": "Vodcasts",
                "Events": "Events",
            }.get(p.category, p.category),
            deck_path="Fortune_Newsletters_2026.pptx",
            slide_number=1,
        )
        for p in products
    ]
    uniq: dict[tuple[str, str], ProductSlideRef] = {}
    for row in rows:
        uniq[(row.product_name, row.category)] = row
    return GtmProductMap(list(uniq.values()))


def _slide_titles(prs: Presentation) -> list[str]:
    return [s.shapes.title.text for s in prs.slides]


def test_build_happy_path_returns_payload():
    generator = _build_generator()
    generator._s3.put_object.return_value = {}
    generator._s3.generate_presigned_url.return_value = "https://s3.example.com/deck.pptx"

    schema = _build_schema()
    source_prs = _presentation_with_n_slides(1)
    product_map = _product_map_for(*schema.confirmed_products)

    with (
        patch("requests.get") as mock_get,
        patch.object(generator, "_load_pptx", return_value=source_prs) as mock_load,
    ):
        mock_get.return_value.content = fortuneai_fixture_bytes()
        mock_get.return_value.raise_for_status = MagicMock()
        result = generator.build(
            schema,
            _FORTUNEAI_URL,
            product_map=product_map,
            audience_data=sample_audience_data(),
            logo_bytes=MINIMAL_PNG,
        )

    assert result["download_url"] == "https://s3.example.com/deck.pptx"
    assert result["client_name"] == "Acme Corp"
    assert result["template_key"] == "FortuneAI_DeckTemplate.pptx"
    assert result["slide_count"] == 10
    assert result["warnings"] == []
    mock_load.assert_called_with("product-decks/Fortune_Newsletters_2026.pptx")
    generator._s3.put_object.assert_called_once()
    put_kwargs = generator._s3.put_object.call_args.kwargs
    assert put_kwargs["Key"].startswith("generated/")
    assert put_kwargs["Key"].endswith(".pptx")


def test_build_rejects_non_fortuneai_template_url():
    generator = _build_generator()
    schema = _build_schema()
    product_map = _product_map_for(*schema.confirmed_products)

    with pytest.raises(ValueError, match="FortuneAI_DeckTemplate"):
        generator.build(
            schema,
            "https://fortune.sharepoint.com/Category_Presentation_Technology.pptx",
            product_map=product_map,
        )


def test_build_missing_gtm_map_raises():
    generator = _build_generator()
    schema = _build_schema()
    product_map = GtmProductMap([])

    with patch("requests.get") as mock_get:
        mock_get.return_value.content = _fortuneai_template_bytes()
        mock_get.return_value.raise_for_status = MagicMock()
        with pytest.raises(ValueError, match="No GTM Product Tags row"):
            generator.build(schema, _FORTUNEAI_URL, product_map=product_map)


def test_build_too_few_template_slides_raises():
    generator = _build_generator()
    schema = _build_schema()
    product_map = _product_map_for(*schema.confirmed_products)

    with patch("requests.get") as mock_get:
        mock_get.return_value.content = _blank_bytes(3)
        mock_get.return_value.raise_for_status = MagicMock()
        with pytest.raises(ValueError, match="at least 19"):
            generator.build(schema, _FORTUNEAI_URL, product_map=product_map)


def test_assemble_omits_empty_dividers_and_keeps_order():
    """Only funded dividers remain; High-Impact before Print; empty Premium omitted."""
    generator = _build_generator()
    schema = _build_schema(
        confirmed_products=[
            Product(
                name="Product A",
                cadence="monthly",
                price=5_000,
                category="Digital Media",
            ),
            Product(
                name="Product B",
                cadence="annual",
                price=35_000,
                category="Print",
            ),
        ]
    )
    source_prs = _presentation_with_n_slides(1)
    source_prs.slides[0].shapes.title.text = "PRODUCT"
    product_map = _product_map_for(*schema.confirmed_products)

    with (
        patch("requests.get") as mock_get,
        patch.object(generator, "_load_pptx", return_value=source_prs),
    ):
        mock_get.return_value.content = _fortuneai_template_bytes(extra_tail=1)
        mock_get.return_value.raise_for_status = MagicMock()
        prs = generator.assemble_skeleton(
            schema, _FORTUNEAI_URL, product_map=product_map
        )

    titles = _slide_titles(prs)
    # 12 spine + DIV-0 + product + DIV-3 + product + investment + thank you
    assert titles[:12] == [f"SPINE-{i}" for i in range(1, 13)]
    assert titles[12] == "DIV-0"  # High-Impact Media
    assert titles[13] == "PRODUCT"
    assert titles[14] == "DIV-3"  # Print (index 3 in divider list)
    assert titles[15] == "PRODUCT"
    assert titles[16] == "INVESTMENT"
    assert titles[17] == "THANK YOU"
    assert "DIV-1" not in titles
    assert "DIV-2" not in titles
    assert "DIV-4" not in titles
    assert "EXTRA-1" not in titles
    assert len(prs.slides) == 18


def test_assemble_single_newsletter_keeps_editorial_divider_only():
    generator = _build_generator()
    schema = _build_schema()
    source_prs = _presentation_with_n_slides(1)
    source_prs.slides[0].shapes.title.text = "NL-PRODUCT"
    product_map = _product_map_for(*schema.confirmed_products)

    with (
        patch("requests.get") as mock_get,
        patch.object(generator, "_load_pptx", return_value=source_prs),
    ):
        mock_get.return_value.content = _fortuneai_template_bytes()
        mock_get.return_value.raise_for_status = MagicMock()
        prs = generator.assemble_skeleton(
            schema, _FORTUNEAI_URL, product_map=product_map
        )

    titles = _slide_titles(prs)
    # 12 spine + Editorial (DIV-1) + product + investment + thank you = 16
    assert len(prs.slides) == 16
    assert titles[12] == "DIV-1"
    assert titles[13] == "NL-PRODUCT"
    assert titles[14] == "INVESTMENT"
    assert titles[15] == "THANK YOU"


def test_assemble_events_product_fails_loud():
    generator = _build_generator()
    schema = _build_schema(
        confirmed_products=[
            Product(
                name="BrainStorm Summit",
                cadence="annual",
                price=100_000,
                category="Events",
            )
        ]
    )
    product_map = _product_map_for(*schema.confirmed_products)

    with patch("requests.get") as mock_get:
        mock_get.return_value.content = _fortuneai_template_bytes()
        mock_get.return_value.raise_for_status = MagicMock()
        with pytest.raises(ValueError, match="escalate"):
            generator.assemble_skeleton(
                schema, _FORTUNEAI_URL, product_map=product_map
            )


def test_assemble_skeleton_clones_exact_slide_number():
    """Exact hit: Deck Path + Slide # from the map, not Titan similarity."""
    generator = _build_generator()
    schema = _build_schema(
        confirmed_products=[
            Product(
                name="CEO Daily",
                cadence="weekly",
                price=20_000,
                category="Newsletter",
            )
        ]
    )
    product_map = GtmProductMap(
        [
            ProductSlideRef(
                product_name="CEO Daily",
                category="Newsletters",
                deck_path="Fortune_Newsletters_2026.pptx",
                slide_number=3,
            )
        ]
    )
    source_prs = _presentation_with_n_slides(3)
    for i, slide in enumerate(source_prs.slides):
        slide.shapes.title.text = f"SRC-{i + 1}"

    with (
        patch("requests.get") as mock_get,
        patch.object(generator, "_load_pptx", return_value=source_prs) as mock_load,
        patch.object(generator, "_clone_slide", wraps=generator._clone_slide) as mock_clone,
    ):
        mock_get.return_value.content = _fortuneai_template_bytes()
        mock_get.return_value.raise_for_status = MagicMock()
        prs = generator.assemble_skeleton(
            schema, _FORTUNEAI_URL, product_map=product_map
        )

    mock_load.assert_called_with("product-decks/Fortune_Newsletters_2026.pptx")
    mock_clone.assert_called_once()
    assert mock_clone.call_args.args[1] == 2  # 0-based index for Slide #3
    assert prs.slides[13].shapes.title.text == "SRC-3"


def test_assemble_skeleton_returns_presentation_without_s3_upload():
    """Seam for Cursor: assemble only — no Anthropic, no upload."""
    generator = _build_generator()
    schema = _build_schema()
    source_prs = _presentation_with_n_slides(1)
    product_map = _product_map_for(*schema.confirmed_products)

    with (
        patch("requests.get") as mock_get,
        patch.object(generator, "_load_pptx", return_value=source_prs),
        patch("anthropic.Anthropic") as mock_anthropic,
    ):
        mock_get.return_value.content = _fortuneai_template_bytes()
        mock_get.return_value.raise_for_status = MagicMock()
        prs = generator.assemble_skeleton(
            schema, _FORTUNEAI_URL, product_map=product_map
        )

    assert len(prs.slides) == 16
    assert hasattr(prs, "slides")
    generator._s3.put_object.assert_not_called()
    generator._s3.generate_presigned_url.assert_not_called()
    mock_anthropic.assert_not_called()


def test_build_loads_fortuneai_from_s3_when_url_omitted():
    generator = _build_generator()
    generator._s3.put_object.return_value = {}
    generator._s3.generate_presigned_url.return_value = "https://s3.example.com/deck.pptx"
    schema = _build_schema()
    product_map = _product_map_for(*schema.confirmed_products)
    source_prs = _presentation_with_n_slides(1)

    def _get_object(**kwargs):
        key = kwargs["Key"]
        body = MagicMock()
        if key == "templates/FortuneAI_DeckTemplate.pptx":
            body.read.return_value = fortuneai_fixture_bytes()
        elif key.endswith(".xlsx"):
            body.read.return_value = b"not-used"
        else:
            buf = io.BytesIO()
            source_prs.save(buf)
            body.read.return_value = buf.getvalue()
        return {"Body": body}

    generator._s3.get_object.side_effect = _get_object

    with patch.object(generator, "_load_pptx", return_value=source_prs) as mock_load:
        result = generator.build(
            schema,
            template_url=None,
            product_map=product_map,
            audience_data=sample_audience_data(),
            logo_bytes=MINIMAL_PNG,
        )

    assert result["template_key"] == "FortuneAI_DeckTemplate.pptx"
    assert result["slide_count"] == 10
    generator._s3.get_object.assert_any_call(
        Bucket="test-bucket", Key="templates/FortuneAI_DeckTemplate.pptx"
    )
    mock_load.assert_called_with("product-decks/Fortune_Newsletters_2026.pptx")


def test_build_delegates_to_assemble_skeleton():
    generator = _build_generator()
    generator._s3.put_object.return_value = {}
    generator._s3.generate_presigned_url.return_value = "https://s3.example.com/deck.pptx"
    schema = _build_schema()
    mock_prs = _presentation_with_n_slides(2)

    with (
        patch.object(generator, "assemble_skeleton", return_value=mock_prs) as mock_assemble,
        patch("ingestion.generator.apply_placeholders", return_value=[]),
    ):
        result = generator.build(
            schema,
            _FORTUNEAI_URL,
            audience_data=sample_audience_data(),
            logo_bytes=MINIMAL_PNG,
        )

    mock_assemble.assert_called_once()
    assert result["slide_count"] == 2
    assert result["download_url"] == "https://s3.example.com/deck.pptx"
    assert result["template_key"] == "FortuneAI_DeckTemplate.pptx"
    generator._s3.put_object.assert_called_once()


def test_build_leaves_product_clone_title_untouched():
    generator = _build_generator()
    generator._s3.put_object.return_value = {}
    generator._s3.generate_presigned_url.return_value = "https://s3.example.com/deck.pptx"
    schema = _build_schema(
        confirmed_products=[
            Product(
                name="CEO Daily",
                cadence="weekly",
                price=50_000,
                category="Newsletter",
            )
        ]
    )
    product_map = GtmProductMap(
        [
            ProductSlideRef(
                product_name="CEO Daily",
                category="Newsletters",
                deck_path="Fortune_Newsletters_2026.pptx",
                slide_number=1,
            )
        ]
    )
    source_prs = _presentation_with_n_slides(1)
    source_prs.slides[0].shapes.title.text = "CEO DAILY"

    with (
        patch("requests.get") as mock_get,
        patch.object(generator, "_load_pptx", return_value=source_prs),
    ):
        mock_get.return_value.content = fortuneai_fixture_bytes()
        mock_get.return_value.raise_for_status = MagicMock()
        result = generator.build(
            schema,
            _FORTUNEAI_URL,
            product_map=product_map,
            audience_data=sample_audience_data(),
            logo_bytes=MINIMAL_PNG,
        )

    body = generator._s3.put_object.call_args.kwargs["Body"]
    built = Presentation(io.BytesIO(body))
    titles = [
        s.shapes.title.text
        for s in built.slides
        if s.shapes.title is not None and s.shapes.title.text
    ]
    assert "CEO DAILY" in titles
    assert result["slide_count"] == 10
