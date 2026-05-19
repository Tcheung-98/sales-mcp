"""  
Step 3 — .pptx parser
  THE PURPOSE OF THIS FILE IS EXCLUSIVELY TO TURN EACH PAGE OF A PPTX INTO A LIST OF SLIDEROW OBJECTS TO BE WRITTEN INTO A DB LATER


  ingestion/parser.py — parse_pptx(bytes) -> list[SlideRecord], content_hash(bytes) -> str. Handles: empty notes → "", missing title → None,
  image-only slides, zero-slide decks, per-slide error isolation.

  prs = Presentation(io.BytesIO(pptx_bytes))  # the whole file
  prs.slides          # iterable of all slides

  slide               # one slide
  slide.shapes        # all shapes on the slide (text boxes, images, charts, etc.)
  slide.shapes.title  # the title placeholder specifically, or None
  slide.slide_layout  # the layout template this slide uses
  slide.slide_layout.name  # e.g. "Title Slide", "Title and Content"
  slide.has_notes_slide   # bool — does this slide have speaker notes?
  slide.notes_slide.notes_text_frame.text  # the raw notes text

  shape               # one shape on a slide
  shape.has_text_frame  # bool — does this shape contain text?
  shape.text_frame    # the text container
  shape.text_frame.paragraphs  # list of paragraphs
  para.text           # the full text of one paragraph
"""

import hashlib
import io
import logging
from pptx import Presentation
from ingestion.models import SlideRow, Tags

logger = logging.getLogger(__name__)

def _clean(text: str) -> str:
    return text.replace('\xa0',' ').replace('\x0b', ' ').strip()


# content_hash to enable nightly diff-sync jobs. exists only for checking diffs for updates to our s3 buckets.
def content_hash(pptx_bytes: bytes) -> str:
    """Generate a content hash for the given pptx bytes."""
    return hashlib.sha256(pptx_bytes).hexdigest()

# the objective:
# this way we can inform the LLM what pages are actually relevant when it comes to generating the slides
def parse_pptx(pptx_bytes, deck_id, source_path, ingested_at, tags) -> list[SlideRow]:
    """Parse the given pptx bytes and return a list of SlideRow objects."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    pptx_hash = content_hash(pptx_bytes)
    res = []
    for i, slide in enumerate(prs.slides, start=1):
        try:
            title_shape = slide.shapes.title
            body_text = [
                _clean(para.text) for shape in slide.shapes if shape.has_text_frame and shape is not title_shape
                for para in shape.text_frame.paragraphs if para.text.strip()
                ]

            cur_slide = SlideRow(
                deck_id=deck_id,
                source_path=source_path,
                content_hash=pptx_hash,
                ingested_at=ingested_at,
                slide_number=i,
                layout_name=slide.slide_layout.name,
                title=_clean(title_shape.text) if title_shape else None,
                body_text=body_text,
                speaker_notes=_clean(slide.notes_slide.notes_text_frame.text) if slide.has_notes_slide else '',
                tags=tags,
                tag_sources=None

            )
            res.append(cur_slide)
        except Exception as e:
            logger.warning("slide %d of deck %s failed to parse: %s", i, deck_id, e)
    return res