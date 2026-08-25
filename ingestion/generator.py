import csv
import io
import json
import logging
import os
import re
from urllib.parse import urlparse
from uuid import uuid4

import anthropic
import boto3
import requests
from docx import Document
from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part, XmlPart
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as PresentationType
from pypdf import PdfReader

from ingestion.audience_data import AudienceData, load_audience_data_from_s3
from ingestion.category_dividers import (
    CATEGORY_DIVIDERS,
    FORTUNEAI_DIVIDER_COUNT,
    FORTUNEAI_FIRST_DIVIDER_INDEX,
    FORTUNEAI_MIN_SLIDES,
    FORTUNEAI_TEMPLATE_BASENAME,
    divider_index_for_category,
    fortuneai_template_key,
    is_fortuneai_template_url,
)
from ingestion.gtm_product_map import (
    GtmProductMap,
    ProductSlideRef,
    load_gtm_product_map_from_s3,
    product_deck_s3_key,
)
from ingestion.placeholder_ai import PlaceholderAI
from ingestion.placeholder_fills import apply_placeholders, fetch_logo_bytes
from ingestion.pptx_tools import (
    apply_replacements,
    delete_slide,
    set_ph_text,
    sync_sections,
)
from ingestion.schema import DeckSchema

logger = logging.getLogger(__name__)

_SECRET_NAME = "fortune-sales-mcp/claude-api-key"
_DEFAULT_MODEL = "claude-sonnet-4-6"
_COST_PER_M_INPUT = 3.00
_COST_PER_M_OUTPUT = 15.00

_TEMPLATE_URL_HOST_SUFFIXES = (".sharepoint.com", ".sharepoint.us", ".microsoft.com")

_RELS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


_SYSTEM_PROMPT_TEMPLATE = """\
You are a senior AE at Fortune Media Group building a custom pitch deck for a specific client.

=== FORTUNE SALES SKILL — RULEBOOK ===
{rulebook_text}
=== END RULEBOOK ===

OUTPUT FORMAT
Return ONLY a JSON array — no prose, no markdown fences. One element per arc slot, in slot order.

Slot 1 (cover) is always:
  {{"action": "cover", "title": "...", "subtitle": "..."}}

All other slots clone a corpus slide:
  {{
    "action": "clone",
    "source_path": "...",
    "slide_number": 7,
    "reasoning": "One sentence: why this slide fits this slot and client.",
    "replacements": {{
      "title": "...",
      "client_name": "Acme Corp",
      "body": ["P1 text", "P2 text", ...]
    }}
  }}

FIELD RULES
- For each non-cover slot, pick exactly one slide from THAT SLOT'S candidates
- Do not use candidates listed under a different slot
- source_path and slide_number must come exactly from the candidates shown for that slot
- reasoning: required on every clone — explain why this slide fits this slot's role
- replacements.client_name: always set to the client name from the brief
- replacements.title: override if needed; omit to keep original
- replacements.eyebrow: override if needed; omit to keep original
- replacements.body: rewrite paragraphs to be client-specific. Match the P1/P2/... structure \
shown. Preserve structural headers; rewrite descriptive copy. 18 words max per bullet. \
OMIT body entirely for: (1) data-heavy slides — audience metrics, bar charts, structured \
tables, When/Where/Who grids, pricing tables; (2) event/conference slides (Brainstorm, \
live events, summits) — their value is in the visual design and event details table, not \
rewritten prose. For investment slides, limit body to 5 bullets maximum — one per product \
line, no sub-bullets, no notes. Only provide body for slides with a clear narrative text \
block and no data visualization or event details table.

SELECTION RULES
- Output exactly one slide per slot, in slot order
- Follow the arc, escalation, and voice rules in the rulebook above
- Use the rate card for accurate product names and pricing\
"""


def _format_slide(s: dict) -> str:
    lines = [
        f"[source: {s.get('source_path', '')} | slide: {s.get('slide_number', '?')}]",
        f"Title: {s.get('title') or '(none)'}",
    ]
    for i, b in enumerate(s.get("body_text") or [], start=1):
        lines.append(f"P{i}: {b}")
    return "\n".join(lines)


class DeckGenerator:
    def __init__(
        self,
        bucket: str | None = None,
        blank_key: str | None = None,
        rulebook_key: str | None = None,
        rate_sheet_key: str | None = None,
        secret_name: str = _SECRET_NAME,
        model: str = _DEFAULT_MODEL,
    ) -> None:
        self._bucket = bucket or os.environ["S3_SNAPSHOT_BUCKET"]
        self._blank_key = blank_key or os.environ.get(
            "PPTX_BLANK_DECK_KEY", "templates/blank.pptx"
        )
        self._rulebook_key = rulebook_key or os.environ.get(
            "RULEBOOK_KEY", "templates/rulebook.docx"
        )
        self._rate_sheet_key = rate_sheet_key or os.environ.get("RATE_SHEET_KEY")
        self._secret_name = secret_name
        self._model = model
        self._s3 = boto3.client("s3")
        self._blank_bytes: bytes | None = None
        self._pptx_cache: dict[str, PresentationType] = {}
        self._rulebook_text: str | None = None
        self._api_key: str | None = None
        self._rate_sheet: str | None = None
        self._gtm_product_map: GtmProductMap | None = None
        self._audience_data: AudienceData | None = None

    @staticmethod
    def _import_ctx(target_prs: PresentationType) -> dict:
        """Per-target bookkeeping for parts pulled in from other packages.

        ``next_partname`` only sees parts already reachable from the package
        root, so track allocations ourselves; ``layouts`` keeps repeat clones of
        one source layout from importing a fresh master each time.
        """
        ctx = getattr(target_prs, "_clone_import_ctx", None)
        if ctx is None:
            ctx = {"package": target_prs.part.package, "taken": set(), "layouts": {}}
            setattr(target_prs, "_clone_import_ctx", ctx)
        ctx["taken"] |= {part.partname for part in ctx["package"].iter_parts()}
        return ctx

    @staticmethod
    def _alloc_partname(ctx: dict, source_partname) -> PackURI:
        stem = source_partname.filename.rsplit(".", 1)[0]
        match = re.match(r"[A-Za-z_]*", stem)
        prefix = (match.group() if match is not None else "") or "part"
        n = 1
        while True:
            candidate = PackURI(
                f"{source_partname.baseURI}/{prefix}{n}.{source_partname.ext}"
            )
            if candidate not in ctx["taken"]:
                ctx["taken"].add(candidate)
                return candidate
            n += 1

    @staticmethod
    def _importable_part(target_prs: PresentationType, source_part):
        """Return a binary part safe to relate into ``target_prs``.

        Product decks and FortuneAI share media partnames (``image63.png``,
        ``hdphoto40.wdp``, …). Relating the source part directly would emit two
        zip entries under one partname, so copy the blob to a free partname
        whenever the name is already taken.
        """
        ctx = DeckGenerator._import_ctx(target_prs)
        if source_part.partname not in ctx["taken"]:
            ctx["taken"].add(source_part.partname)
            return source_part
        return Part(
            DeckGenerator._alloc_partname(ctx, source_part.partname),
            source_part.content_type,
            ctx["package"],
            source_part.blob,
        )

    @staticmethod
    def _copy_xml_part(ctx: dict, source_part) -> XmlPart:
        # Keep the source's own part class (SlideLayoutPart, SlideMasterPart, …)
        # so the copy still answers python-pptx's object-model traversals. Parts
        # python-pptx has no class for — themes — load blob-backed, so read them
        # as generic XML instead.
        source_cls = type(source_part)
        part_cls = source_cls if issubclass(source_cls, XmlPart) else XmlPart
        return part_cls(
            DeckGenerator._alloc_partname(ctx, source_part.partname),
            source_part.content_type,
            ctx["package"],
            parse_xml(source_part.blob),
        )

    @staticmethod
    def _import_slide_layout(target_prs: PresentationType, source_layout):
        """Copy ``source_layout`` — plus its master and theme — into ``target_prs``.

        Cloned product slides keep placeholders, and placeholders inherit font,
        size, colour and bullet formatting from their layout/master/theme. Point
        them at a FortuneAI layout instead and the slide silently restyles, so
        bring the real chain along. The copied master is pruned to just this
        layout, which keeps the import to roughly the layout's own weight.
        """
        ctx = DeckGenerator._import_ctx(target_prs)
        cache_key = id(source_layout.part)
        if cache_key in ctx["layouts"]:
            return ctx["layouts"][cache_key]

        src_layout_part = source_layout.part
        src_master_part = source_layout.slide_master.part

        master_part = DeckGenerator._copy_xml_part(ctx, src_master_part)
        master_map: dict[str, str] = {}
        for rId, rel in src_master_part.rels.items():
            # Layout list is rebuilt below so the copy carries only what we need.
            if rel.reltype == RT.SLIDE_LAYOUT:
                continue
            if rel.is_external:
                master_map[rId] = master_part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            elif rel.reltype == RT.THEME:
                master_map[rId] = master_part.relate_to(
                    DeckGenerator._import_theme(target_prs, ctx, rel.target_part),
                    rel.reltype,
                )
            else:
                master_map[rId] = master_part.relate_to(
                    DeckGenerator._importable_part(target_prs, rel.target_part),
                    rel.reltype,
                )

        layout_part = DeckGenerator._copy_xml_part(ctx, src_layout_part)
        layout_map: dict[str, str] = {}
        for rId, rel in src_layout_part.rels.items():
            if rel.reltype == RT.SLIDE_MASTER:
                layout_map[rId] = layout_part.relate_to(master_part, rel.reltype)
            elif rel.is_external:
                layout_map[rId] = layout_part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            else:
                layout_map[rId] = layout_part.relate_to(
                    DeckGenerator._importable_part(target_prs, rel.target_part),
                    rel.reltype,
                )
        DeckGenerator._remap_rids(layout_part._element, layout_map)

        DeckGenerator._remap_rids(master_part._element, master_map)
        layout_id_lst = master_part._element.find(qn("p:sldLayoutIdLst"))
        if layout_id_lst is None:
            raise ValueError(
                f"slide master {src_master_part.partname} has no sldLayoutIdLst"
            )
        for child in list(layout_id_lst):
            layout_id_lst.remove(child)
        layout_id_lst.append(
            parse_xml(
                f'<p:sldLayoutId xmlns:p="{_PML_NS}" xmlns:r="{_RELS_NS}" '
                f'id="2147483649" r:id="{master_part.relate_to(layout_part, RT.SLIDE_LAYOUT)}"/>'
            )
        )

        master_id_lst = target_prs.part._element.find(qn("p:sldMasterIdLst"))
        existing = [int(e.get("id")) for e in master_id_lst if e.get("id")]
        master_id_lst.append(
            parse_xml(
                f'<p:sldMasterId xmlns:p="{_PML_NS}" xmlns:r="{_RELS_NS}" '
                f'id="{max(existing, default=2147483648) + 1}" '
                f'r:id="{target_prs.part.relate_to(master_part, RT.SLIDE_MASTER)}"/>'
            )
        )

        ctx["layouts"][cache_key] = layout_part
        return layout_part

    @staticmethod
    def _import_theme(target_prs: PresentationType, ctx: dict, source_theme_part):
        theme_part = DeckGenerator._copy_xml_part(ctx, source_theme_part)
        theme_map: dict[str, str] = {}
        for rId, rel in source_theme_part.rels.items():
            if rel.is_external:
                theme_map[rId] = theme_part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            else:
                theme_map[rId] = theme_part.relate_to(
                    DeckGenerator._importable_part(target_prs, rel.target_part),
                    rel.reltype,
                )
        DeckGenerator._remap_rids(theme_part._element, theme_map)
        return theme_part

    @staticmethod
    def _remap_rids(element, rId_map: dict[str, str]) -> None:
        """Rewrite every relationship-namespace attribute (r:embed, r:id, r:link, …).

        Single pass over attributes: sequential string replacement would corrupt
        the mapping whenever a new rId collides with an old one still unwritten.
        """
        prefix = f"{{{_RELS_NS}}}"
        for el in element.iter():
            for attr in list(el.keys()):
                if not attr.startswith(prefix):
                    continue
                new_rId = rId_map.get(el.get(attr))
                if new_rId is not None:
                    el.set(attr, new_rId)

    @staticmethod
    def _clone_slide(
        source_prs: PresentationType, slide_idx: int, target_prs: PresentationType
    ):
        source_slide = source_prs.slides[slide_idx]

        # Add a placeholder slide — gives us a proper slide part + sldIdLst entry
        new_slide = target_prs.slides.add_slide(target_prs.slide_layouts[0])

        # Carry every source relationship except the layout (rewired below) and
        # notes (deliberately dropped). Copying only images leaves hyperlinks,
        # hdphoto sidecars and media as dangling r:ids, which makes PowerPoint
        # declare the deck corrupt and "repair" it on open.
        rId_map: dict[str, str] = {}
        for rId, rel in source_slide.part.rels.items():
            if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
                continue
            if rel.is_external:
                new_rId = new_slide.part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            else:
                new_rId = new_slide.part.relate_to(
                    DeckGenerator._importable_part(
                        target_prs, source_slide.part.related_part(rId)
                    ),
                    rel.reltype,
                )
            rId_map[rId] = new_rId

        # Serialize source cSld and re-parse with pptx element classes (parse_xml,
        # not etree.fromstring, so spTree and other pptx attrs are available)
        fixed_cSld = parse_xml(etree.tostring(source_slide._element.cSld))
        DeckGenerator._remap_rids(fixed_cSld, rId_map)

        # Swap target slide's cSld with the fixed clone
        tgt_sld = new_slide._element
        tgt_sld.replace(tgt_sld.cSld, fixed_cSld)

        # shapes is a lazyproperty cached during add_slide; invalidate so next access
        # gets a fresh SlideShapes pointing at the new spTree, not the discarded one
        new_slide.__dict__.pop("shapes", None)

        # Point the clone at the source slide's own layout, imported into this
        # package. Matching a FortuneAI layout by name instead would leave the
        # slide's placeholders inheriting the wrong fonts, bullets and colours.
        layout_part = DeckGenerator._import_slide_layout(
            target_prs, source_slide.slide_layout
        )
        for rId, rel in list(new_slide.part.rels.items()):
            if rel.reltype == RT.SLIDE_LAYOUT:
                new_slide.part.drop_rel(rId)
                break
        new_slide.part.relate_to(layout_part, RT.SLIDE_LAYOUT)

        return new_slide

    @staticmethod
    def _delete_slide(prs: PresentationType, slide_idx: int) -> None:
        delete_slide(prs, slide_idx)

    @staticmethod
    def _insert_slide_at(prs: PresentationType, position: int) -> None:
        # _clone_slide always appends; move the last sldId to the target position
        sld_id_lst = prs.slides._sldIdLst
        sld_id = sld_id_lst[-1]
        sld_id_lst.remove(sld_id)
        sld_id_lst.insert(position, sld_id)

    @staticmethod
    def _validate_template_url(template_url: str) -> None:
        parsed = urlparse(template_url)
        if parsed.scheme != "https":
            raise ValueError("template_url must use HTTPS")
        host = (parsed.hostname or "").lower()
        if not host:
            raise ValueError("template_url must include a host")
        extra_hosts = {
            h.strip().lower()
            for h in os.environ.get("TEMPLATE_URL_ALLOWED_HOSTS", "").split(",")
            if h.strip()
        }
        if host in extra_hosts or any(
            host.endswith(suffix) for suffix in _TEMPLATE_URL_HOST_SUFFIXES
        ):
            return
        raise ValueError(f"template_url host not allowed: {host}")

    def _load_pptx(self, s3_key: str) -> PresentationType:
        if s3_key not in self._pptx_cache:
            resp = self._s3.get_object(Bucket=self._bucket, Key=s3_key)
            self._pptx_cache[s3_key] = Presentation(io.BytesIO(resp["Body"].read()))
            logger.info("loaded pptx from s3://%s/%s", self._bucket, s3_key)
        return self._pptx_cache[s3_key]

    def _load_rulebook(self) -> str:
        if self._rulebook_text is None:
            resp = self._s3.get_object(Bucket=self._bucket, Key=self._rulebook_key)
            body = resp["Body"].read()
            if self._rulebook_key.endswith(".pdf"):
                pdf = PdfReader(io.BytesIO(body))
                self._rulebook_text = "\n".join(
                    page.extract_text() for page in pdf.pages if page.extract_text()
                )
            else:
                doc = Document(io.BytesIO(body))
                self._rulebook_text = "\n".join(
                    p.text for p in doc.paragraphs if p.text.strip()
                )
            logger.info("loaded rulebook from s3://%s/%s", self._bucket, self._rulebook_key)
        return self._rulebook_text

    def _get_api_key(self) -> str:
        if self._api_key is None:
            # Local dev: ANTHROPIC_API_KEY in .env bypasses Secrets Manager
            if env_key := os.environ.get("ANTHROPIC_API_KEY"):
                self._api_key = env_key
                return env_key
            sm = boto3.client("secretsmanager", region_name="us-east-1")
            resp = sm.get_secret_value(SecretId=self._secret_name)
            raw = resp["SecretString"]
            # Handle JSON-wrapped secrets e.g. {"api_key": "sk-ant-..."}
            try:
                parsed = json.loads(raw)
                self._api_key = next(iter(parsed.values())) if isinstance(parsed, dict) else raw
            except (json.JSONDecodeError, StopIteration):
                self._api_key = raw
        return self._api_key

    def _call_claude(
        self, brief: str, arc: list[dict], context_by_slot: dict[int, list[dict]]
    ) -> list[dict]:
        context_parts = []
        for slot in arc:
            if not slot.get("query"):
                continue
            candidates = context_by_slot.get(slot["slot"], [])
            if not candidates:
                continue
            context_parts.append(
                f"=== Slot {slot['slot']}: {slot['role'].upper()} ===\n"
                + "\n\n".join(_format_slide(s) for s in candidates)
            )
        context_text = "\n\n".join(context_parts)

        arc_summary = "\n".join(
            f"  Slot {s['slot']}: {s['role']}" for s in arc
        )

        if self._rate_sheet is None and self._rate_sheet_key:
            resp = self._s3.get_object(Bucket=self._bucket, Key=self._rate_sheet_key)
            reader = csv.DictReader(io.StringIO(resp["Body"].read().decode("cp1252")))
            products = [row for row in reader if row.get("Inventory", "").strip()]
            lines = []
            for p in products:
                parts = [p.get("Inventory", "").strip(), p.get("Product Category", "").strip()]
                if p.get("Cadence", "").strip():
                    parts.append(p["Cadence"].strip())
                if p.get("Vertical", "").strip():
                    parts.append(f"Verticals: {p['Vertical'].strip()}")
                if p.get("Audience Alignment", "").strip():
                    parts.append(f"Audience: {p['Audience Alignment'].strip()}")
                if p.get("Contextual Alignment", "").strip():
                    parts.append(f"Context: {p['Contextual Alignment'].strip()}")
                pricing = []
                for label, col in [
                    ("Daily", "Daily_Cost"), ("Weekly", "Weekly_Cost"),
                    ("Monthly", "Monthly_Cost"), ("Quarterly", "Quarterly_Cost"),
                    ("Half-Year", "Half_Year_Cost"), ("Annual", "Annual_Cost"),
                    ("CPM", "CPM_Rate"), ("Min", "Product_Minimum"), ("Flat Fee", "Flat_Fee"),
                ]:
                    val = p.get(col, "").strip()
                    if val:
                        prefix = "" if val.startswith("$") else "$"
                        pricing.append(f"{label}: {prefix}{val}")
                if pricing:
                    parts.append(" | ".join(pricing))
                if p.get("Notes", "").strip():
                    parts.append(f"Notes: {p['Notes'].strip()}")
                lines.append(" | ".join(filter(None, parts)))
            self._rate_sheet = "\n".join(lines)
            logger.info("loaded rate sheet: %d products", len(products))
        rate_section = (
            f"\n\nFortune product & rate card (use for accurate product names and pricing):\n"
            f"{self._rate_sheet}"
            if self._rate_sheet else ""
        )

        user_msg = (
            f"Brief:\n{brief}"
            f"\n\nDeck arc to fill:\n{arc_summary}"
            f"{rate_section}"
            f"\n\nCandidate slides organized by slot:\n{context_text}"
        )

        system_prompt = _SYSTEM_PROMPT_TEMPLATE.format(
            rulebook_text=self._load_rulebook()
        )

        client = anthropic.Anthropic(api_key=self._get_api_key())
        response = client.messages.create(
            model=self._model,
            max_tokens=4096,
            system=system_prompt,
            messages=[{"role": "user", "content": user_msg}],
        )

        logger.info(
            "claude usage — input: %d tokens, output: %d tokens, cost: $%.4f",
            response.usage.input_tokens,
            response.usage.output_tokens,
            (response.usage.input_tokens / 1_000_000 * _COST_PER_M_INPUT)
            + (response.usage.output_tokens / 1_000_000 * _COST_PER_M_OUTPUT),
        )
        block = response.content[0]
        if not isinstance(block, anthropic.types.TextBlock):
            raise ValueError(f"Unexpected response block type: {type(block)}")
        raw = block.text.strip()
        # Strip markdown fences if Claude ignored the prompt instruction
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[-1]
            raw = raw.rsplit("```", 1)[0].strip()
        logger.debug("claude raw response: %s", raw[:200])
        slides = json.loads(raw)
        if not isinstance(slides, list):
            raise ValueError(f"Claude returned non-list JSON: {type(slides)}")
        return slides

    @staticmethod
    def _dedup_shape_ids(slide, output_prs) -> None:
        """Renumber cNvPr ids on a freshly-cloned slide to avoid collisions across slides."""
        used: set[int] = set()
        for s in output_prs.slides:
            if s._element is slide._element:
                continue
            for el in s._element.iter(qn("p:cNvPr")):
                try:
                    used.add(int(el.get("id", 0)))
                except (ValueError, TypeError):
                    pass

        def _sid(el) -> int:
            try:
                return int(el.get("id", 0))
            except (ValueError, TypeError):
                return 0

        own = list(slide._element.iter(qn("p:cNvPr")))
        # Reserve ids further down this slide too, otherwise a replacement id can
        # collide with a shape we have not visited yet and duplicate it.
        reserved = used | {_sid(el) for el in own}
        next_id = max(reserved, default=0) + 1
        taken: set[int] = set()
        for el in own:
            sid = _sid(el)
            if sid == 0 or sid in used or sid in taken:
                while next_id in reserved:
                    next_id += 1
                el.set("id", str(next_id))
                reserved.add(next_id)
                taken.add(next_id)
                next_id += 1
            else:
                taken.add(sid)

    @staticmethod
    def _set_ph_text(ph, text: str) -> None:
        set_ph_text(ph, text)

    @staticmethod
    def _apply_replacements(slide, replacements: dict) -> None:
        apply_replacements(slide, replacements)

    def _build_pptx(self, slides: list[dict]) -> bytes:
        if self._blank_bytes is None:
            resp = self._s3.get_object(Bucket=self._bucket, Key=self._blank_key)
            self._blank_bytes = resp["Body"].read()
            logger.info("loaded blank template from s3://%s/%s", self._bucket, self._blank_key)
        blank_prs = Presentation(io.BytesIO(self._blank_bytes))  # cover slide source
        output_prs = Presentation(io.BytesIO(self._blank_bytes))  # output deck

        # Drop relationship AND sldId element for each seed slide
        for sld_id in list(output_prs.slides._sldIdLst):
            r_id = sld_id.get(qn("r:id"))
            output_prs.part.drop_rel(r_id)
            output_prs.slides._sldIdLst.remove(sld_id)

        for slide_data in slides:
            action = slide_data.get("action", "clone")
            if action == "cover":
                slide = self._clone_slide(blank_prs, 0, output_prs)
                DeckGenerator._dedup_shape_ids(slide, output_prs)
                ph_map = {
                    ph.placeholder_format.idx: ph
                    for ph in slide.placeholders
                    if ph.has_text_frame
                }
                for title_idx in (12, 0):
                    if ph := ph_map.get(title_idx):
                        DeckGenerator._set_ph_text(ph, slide_data.get("title", ""))
                        break
                for sub_idx in (10, 1):
                    if ph := ph_map.get(sub_idx):
                        DeckGenerator._set_ph_text(ph, slide_data.get("subtitle", ""))
                        break
            elif action == "clone":
                source_prs = self._load_pptx(slide_data["source_path"])
                slide = self._clone_slide(source_prs, slide_data["slide_number"] - 1, output_prs)
                DeckGenerator._dedup_shape_ids(slide, output_prs)
                self._apply_replacements(slide, slide_data.get("replacements", {}))
            else:
                raise ValueError(f"Unknown slide action: {action!r}")

        buf = io.BytesIO()
        output_prs.save(buf)
        return buf.getvalue()

    def _get_gtm_product_map(self) -> GtmProductMap:
        if self._gtm_product_map is None:
            self._gtm_product_map = load_gtm_product_map_from_s3(self._s3, self._bucket)
        return self._gtm_product_map

    def _get_audience_data(self) -> AudienceData:
        if self._audience_data is None:
            self._audience_data = load_audience_data_from_s3(self._s3, self._bucket)
        return self._audience_data

    def _load_fortuneai_template(
        self, template_url: str | None
    ) -> tuple[PresentationType, str]:
        """Load FortuneAI_DeckTemplate from SharePoint URL or S3.

        Returns (presentation, template_key used in build payload).
        Always returns a fresh Presentation — assembly mutates the deck.
        """
        if template_url:
            self._validate_template_url(template_url)
            if not is_fortuneai_template_url(template_url):
                raise ValueError(
                    "template_url must point to FortuneAI_DeckTemplate "
                    f"(got {template_url.split('?', 1)[0]!r})"
                )
            resp = requests.get(template_url, timeout=30)
            resp.raise_for_status()
            return Presentation(io.BytesIO(resp.content)), FORTUNEAI_TEMPLATE_BASENAME

        key = fortuneai_template_key()
        try:
            resp = self._s3.get_object(Bucket=self._bucket, Key=key)
            data = resp["Body"].read()
        except Exception as exc:
            raise ValueError(
                f"Failed to load FortuneAI template from "
                f"s3://{self._bucket}/{key}: {exc}"
            ) from exc
        logger.info("loaded FortuneAI template from s3://%s/%s", self._bucket, key)
        return Presentation(io.BytesIO(data)), FORTUNEAI_TEMPLATE_BASENAME

    def _group_products_by_divider(
        self,
        schema: DeckSchema,
        gtm_map: GtmProductMap,
    ) -> list[list[ProductSlideRef]]:
        """Bucket confirmed products into the five Workflow dividers (A5 lookup)."""
        groups: list[list[ProductSlideRef]] = [[] for _ in CATEGORY_DIVIDERS]
        failures: list[str] = []
        for product in schema.confirmed_products:
            try:
                divider_i = divider_index_for_category(product.category)
                ref = gtm_map.lookup(product.name, product.category)
            except ValueError as exc:
                failures.append(str(exc))
                continue
            groups[divider_i].append(ref)
        if failures:
            raise ValueError(
                "FortuneAI product placement failed: " + "; ".join(failures)
            )
        return groups

    def _clone_product_ref(
        self, ref: ProductSlideRef, target_prs: PresentationType
    ):
        s3_key = product_deck_s3_key(ref.deck_path)
        try:
            source_prs = self._load_pptx(s3_key)
        except Exception as exc:
            raise ValueError(
                f"Failed to load Deck Path {ref.deck_path!r} for product "
                f"{ref.product_name!r} (s3://{self._bucket}/{s3_key}): {exc}"
            ) from exc
        if ref.slide_number > len(source_prs.slides):
            raise ValueError(
                f"Slide # {ref.slide_number} out of range for product "
                f"{ref.product_name!r} in {ref.deck_path!r} "
                f"({len(source_prs.slides)} slides)"
            )
        new_slide = self._clone_slide(source_prs, ref.slide_number - 1, target_prs)
        self._dedup_shape_ids(new_slide, target_prs)
        return new_slide

    def assemble_skeleton(
        self,
        schema: DeckSchema,
        template_url: str | None = None,
        product_map: GtmProductMap | None = None,
    ) -> PresentationType:
        """Assemble FortuneAI spine + conditional dividers + exact GTM product clones.

        Keeps intro/narrative/investment/thank-you stock (C2 fills in ``build``). Inserts
        category dividers only when ≥1 funded product maps to that section. Product
        pages are wholesale A5 clones — no AI edits.
        """
        gtm_map = product_map if product_map is not None else self._get_gtm_product_map()
        # Resolve placement before mutating the template so map misses fail loud early.
        groups = self._group_products_by_divider(schema, gtm_map)
        prs, _template_key = self._load_fortuneai_template(template_url)

        if len(prs.slides) < FORTUNEAI_MIN_SLIDES:
            raise ValueError(
                f"FortuneAI_DeckTemplate must have at least {FORTUNEAI_MIN_SLIDES} "
                f"slides (intro/narrative + 5 dividers + investment + thank you); "
                f"got {len(prs.slides)}"
            )

        # Drop trailing extras (e.g. stock Print slide 20) so post-spine is
        # investment + thank you only before we insert product clones.
        while len(prs.slides) > FORTUNEAI_MIN_SLIDES:
            self._delete_slide(prs, len(prs.slides) - 1)

        first = FORTUNEAI_FIRST_DIVIDER_INDEX
        # Delete unfunded dividers back-to-front so earlier indices stay stable.
        for i in range(FORTUNEAI_DIVIDER_COUNT - 1, -1, -1):
            if not groups[i]:
                self._delete_slide(prs, first + i)

        # Remaining dividers are packed starting at `first`. Insert A5 clones
        # after each kept divider.
        cursor = first
        for refs in groups:
            if not refs:
                continue
            for j, ref in enumerate(refs):
                self._clone_product_ref(ref, prs)
                self._insert_slide_at(prs, cursor + 1 + j)
            cursor += 1 + len(refs)

        sync_sections(prs)
        return prs

    def build(
        self,
        schema: DeckSchema,
        template_url: str | None = None,
        product_map: GtmProductMap | None = None,
        audience_data: AudienceData | None = None,
        logo_bytes: bytes | None = None,
    ) -> dict:
        """Assemble FortuneAI PPTX, fill placeholders, upload.

        template_url: optional pre-authenticated SharePoint download URL for
        FortuneAI_DeckTemplate. When omitted, loads from S3 (FORTUNEAI_TEMPLATE_KEY).
        Product pages use exact GTM Deck Path / Slide #. C2 fills date, logo,
        history, audience Reach/Index, program types, investment, and bounded
        Claude copy for intro, Opportunity, audience title, and program blurbs.
        No stylist.
        """
        prs = self.assemble_skeleton(schema, template_url, product_map=product_map)
        audience = (
            audience_data if audience_data is not None else self._get_audience_data()
        )
        logo = (
            logo_bytes if logo_bytes is not None else fetch_logo_bytes(schema.client_logo)
        )
        ai = PlaceholderAI.from_anthropic(
            anthropic.Anthropic(api_key=self._get_api_key()),
            model=self._model,
        )
        warnings = apply_placeholders(
            prs, schema, audience=audience, logo_bytes=logo, ai=ai
        )
        buf = io.BytesIO()
        prs.save(buf)
        key = f"generated/{uuid4()}.pptx"
        self._s3.put_object(Bucket=self._bucket, Key=key, Body=buf.getvalue())
        url = self._s3.generate_presigned_url(
            "get_object",
            Params={"Bucket": self._bucket, "Key": key},
            ExpiresIn=86400,
        )
        logger.info("deck uploaded to s3://%s/%s", self._bucket, key)
        return {
            "download_url": url,
            "slide_count": len(prs.slides),
            "client_name": schema.client_name,
            "template_key": FORTUNEAI_TEMPLATE_BASENAME,
            "warnings": warnings,
        }

