"""Guardrailed PPTX helpers for skeleton assembly and placeholder fills.

Clone/delete of slides stay on DeckGenerator; this module owns text-apply,
logo insert, and investment-box clone so fills can run without importing
the full generator.
"""

from __future__ import annotations

import io
import re
from copy import deepcopy

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

_P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"

# FortuneAI History of Trust / Audience possessive tokens use U+2019.
APOS = "\u2019"
CLIENT_NAME_TOKEN = "[client name]"
CLIENT_NAME_POSSESSIVE_TOKEN = f"[CLIENT NAME{APOS}S]"
LOGO_TOKEN = "[LOGO]"

_CLIENT_NAME_PLACEHOLDERS = re.compile(
    r"your brand|your company|client name|your organization",
    re.IGNORECASE,
)

_PNG_MAGIC = b"\x89PNG\r\n\x1a\n"
_JPEG_MAGIC = b"\xff\xd8"
_GIF_MAGIC = b"GIF8"


def set_ph_text(ph, text: str) -> None:
    """Replace placeholder text while preserving run-level formatting."""
    tf = ph.text_frame
    first_para = tf.paragraphs[0]
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    runs = first_para.runs
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        first_para.text = text


def apply_replacements(slide, replacements: dict) -> None:
    """Apply title/eyebrow/body/client_name replacements on a slide."""
    if not replacements:
        return

    title = replacements.get("title")
    eyebrow = replacements.get("eyebrow")
    client_name = replacements.get("client_name")
    body = replacements.get("body")

    ph_map = {
        ph.placeholder_format.idx: ph
        for ph in slide.placeholders
        if ph.has_text_frame
    }

    if title and (ph := ph_map.get(0)):
        set_ph_text(ph, title)
    if eyebrow and (ph := ph_map.get(19)):
        set_ph_text(ph, eyebrow)

    # body: rewrite paragraphs one-for-one across non-title/non-eyebrow shapes.
    # body_text[] in the retriever is scraped in slide order, so traverse the same way.
    if body:
        _HANDLED_IDX = {0, 19}
        body_shapes = []
        for s in slide.shapes:
            if not s.has_text_frame:
                continue
            try:
                idx = s.placeholder_format.idx
            except ValueError:
                continue
            if idx in _HANDLED_IDX:
                continue
            body_shapes.append(s)
        all_paras = [p for s in body_shapes for p in s.text_frame.paragraphs]
        for i, new_text in enumerate(body):
            if i >= len(all_paras):
                break
            para = all_paras[i]
            runs = para.runs
            if runs:
                runs[0].text = new_text
                for r in runs[1:]:
                    r._r.getparent().remove(r._r)
            else:
                para.text = new_text

    if client_name:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if _CLIENT_NAME_PLACEHOLDERS.search(run.text):
                        run.text = _CLIENT_NAME_PLACEHOLDERS.sub(
                            client_name, run.text
                        )


def iter_shapes(slide):
    """Yield every shape on a slide, including those nested in groups."""

    def _walk(container):
        for shape in container.shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _walk(shape)

    yield from _walk(slide)


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(p.text or "" for p in shape.text_frame.paragraphs)


def _replace_token_in_paragraph(para, token: str, text: str, *, limit: int | None) -> int:
    """Replace ``token`` in one paragraph. ``limit`` is max replacements (None = all)."""
    remaining = limit
    hits = 0
    run_hit = False
    for run in para.runs:
        if remaining is not None and remaining <= 0:
            return hits
        if token not in run.text:
            continue
        n = run.text.count(token)
        if remaining is not None:
            n = min(n, remaining)
            run.text = run.text.replace(token, text, n)
            remaining -= n
        else:
            run.text = run.text.replace(token, text)
        hits += n
        run_hit = True
    if run_hit:
        return hits
    full = para.text or ""
    if token not in full:
        return 0
    n = full.count(token)
    if remaining is not None:
        n = min(n, remaining)
        new = full.replace(token, text, n)
    else:
        new = full.replace(token, text)
    if para.runs:
        para.runs[0].text = new
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.text = new
    return n


def replace_token(slide, token: str, text: str) -> int:
    """Replace exact ``token`` substrings in all text frames. Returns hit count.

    Prefers in-run replace so formatting stays. If the token is split across
    runs in a paragraph, the paragraph is collapsed onto the first run.
    """
    if not token:
        raise ValueError("token must be non-empty")

    hits = 0
    for shape in iter_shapes(slide):
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            hits += _replace_token_in_paragraph(para, token, text, limit=None)
    return hits


def replace_first_token(slide, token: str, text: str) -> int:
    """Replace the first ``token`` on the slide (document order). Returns 0 or 1."""
    if not token:
        raise ValueError("token must be non-empty")
    for shape in iter_shapes(slide):
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            n = _replace_token_in_paragraph(para, token, text, limit=1)
            if n:
                return n
    return 0


def clone_shape_below(slide, shape, extra_top_emu: int):
    """Deep-copy ``shape`` on the same slide, shifted down by ``extra_top_emu``.

    Used to repeat the Investment category text box once per funded divider.
    Fails loud if the shape has no position transform to copy.
    """
    if extra_top_emu <= 0:
        raise ValueError("Investment category box clone failed: offset must be positive")
    src = shape._element
    new_el = deepcopy(src)
    used: set[int] = set()
    for el in slide._element.iter(qn("p:cNvPr")):
        try:
            used.add(int(el.get("id", 0)))
        except (TypeError, ValueError):
            pass
    nv = new_el.find(qn("p:nvSpPr"))
    cnv = None if nv is None else nv.find(qn("p:cNvPr"))
    if cnv is None:
        raise ValueError("Investment category box clone failed: missing shape id")
    cnv.set("id", str(max(used, default=0) + 1))

    sp_pr = new_el.find(qn("p:spPr"))
    xfrm = None if sp_pr is None else sp_pr.find(qn("a:xfrm"))
    off = None if xfrm is None else xfrm.find(qn("a:off"))
    if off is None or off.get("y") is None:
        raise ValueError("Investment category box clone failed: missing position")
    off.set("y", str(int(off.get("y")) + extra_top_emu))

    slide._element.cSld.spTree.append(new_el)
    slide.__dict__.pop("shapes", None)
    for candidate in slide.shapes:
        if candidate._element is new_el:
            return candidate
    raise ValueError("Investment category box clone failed: clone not on slide")


def _is_picture_placeholder(shape) -> bool:
    try:
        return shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    except (ValueError, AttributeError):
        return False


def _find_logo_placeholder(slide):
    pictures = [s for s in iter_shapes(slide) if _is_picture_placeholder(s)]
    labeled = [s for s in pictures if LOGO_TOKEN in _shape_text(s)]
    if len(labeled) == 1:
        return labeled[0]
    if len(labeled) > 1:
        raise ValueError("slide has more than one [LOGO] picture placeholder")
    idx11 = []
    for shape in pictures:
        try:
            if shape.placeholder_format.idx == 11:
                idx11.append(shape)
        except (ValueError, AttributeError):
            continue
    if len(idx11) == 1:
        return idx11[0]
    if len(pictures) == 1:
        return pictures[0]
    if pictures:
        raise ValueError(
            "slide has multiple picture placeholders and none are [LOGO] "
            "or placeholder idx 11"
        )
    return None


def _validate_logo_bytes(image_bytes: bytes) -> None:
    if not image_bytes:
        raise ValueError("logo image is empty")
    if image_bytes.startswith(_PNG_MAGIC):
        return
    if image_bytes.startswith(_JPEG_MAGIC):
        return
    if image_bytes.startswith(_GIF_MAGIC):
        return
    raise ValueError("logo image is unreadable (expected PNG, JPEG, or GIF)")


def insert_logo(slide, image_bytes: bytes) -> None:
    """Insert ``image_bytes`` into the slide's [LOGO] picture placeholder.

    FortuneAI intro/thanks use a PICTURE placeholder (idx 11) whose text is
    ``[LOGO]``. Fixtures may use any single picture placeholder. Fails loud if
    the placeholder is missing or the bytes are not a PNG/JPEG/GIF.
    """
    _validate_logo_bytes(image_bytes)
    placeholder = _find_logo_placeholder(slide)
    if placeholder is None:
        raise ValueError("logo placeholder not found ([LOGO] picture placeholder)")
    try:
        placeholder.insert_picture(io.BytesIO(image_bytes))
    except Exception as exc:
        raise ValueError(f"logo image is unreadable: {exc}") from exc


def _section_lst(prs):
    """PowerPoint's section list, stored as an extension on p:presentation."""
    ext_lst = prs.part._element.find(qn("p:extLst"))
    if ext_lst is None:
        return None
    for ext in ext_lst:
        section_lst = ext.find(f"{{{_P14_NS}}}sectionLst")
        if section_lst is not None:
            return section_lst
    return None


def _prune_empty_sections(section_lst) -> None:
    for section in list(section_lst):
        id_lst = section.find(f"{{{_P14_NS}}}sldIdLst")
        if id_lst is None or len(id_lst) == 0:
            section_lst.remove(section)
    if len(section_lst) == 0:
        ext = section_lst.getparent()
        ext.getparent().remove(ext)


def delete_slide(prs, slide_idx: int) -> None:
    """Remove a slide by 0-based index (sldIdLst + relationship + section entry)."""
    n = len(prs.slides)
    if slide_idx < 0 or slide_idx >= n:
        raise ValueError(f"slide index {slide_idx} out of range (0..{n - 1})")
    sld_id_lst = prs.slides._sldIdLst
    sld_id = sld_id_lst[slide_idx]
    r_id = sld_id.get(qn("r:id"))
    slide_id = sld_id.get("id")
    prs.part.drop_rel(r_id)
    sld_id_lst.remove(sld_id)

    # A section entry outliving its slide is unreadable content to PowerPoint,
    # which repairs the deck on open. Deletions run from several call sites
    # (unfunded dividers, stock tail, unused audience pages), so keep the
    # section list correct here rather than at each of them.
    section_lst = _section_lst(prs)
    if section_lst is None:
        return
    for section in section_lst:
        id_lst = section.find(f"{{{_P14_NS}}}sldIdLst")
        if id_lst is None:
            continue
        for entry in list(id_lst):
            if entry.get("id") == slide_id:
                id_lst.remove(entry)
    _prune_empty_sections(section_lst)


def sync_sections(prs) -> None:
    """Reconcile the section list with the slides that actually exist.

    FortuneAI ships sections (TITLE, AUDIENCE, PRODUCT OFFERINGS, …) and
    assembly deletes unfunded dividers and the stock Print page, then inserts
    product clones. A section entry left pointing at a deleted slide is content
    PowerPoint cannot read, so it repairs the deck on open. Slides added since
    join the section of the nearest slide before them, which puts product clones
    under the divider they were inserted after.
    """
    section_lst = _section_lst(prs)
    if section_lst is None:
        return

    order: list[str] = []
    for sld_id in prs.slides._sldIdLst:
        slide_id = sld_id.get("id")
        if slide_id is not None:
            order.append(slide_id)
    live = set(order)

    sections = []
    for section in list(section_lst):
        id_lst = section.find(f"{{{_P14_NS}}}sldIdLst")
        if id_lst is None:
            section_lst.remove(section)
            continue
        sections.append((section, id_lst))

    owner: dict[str, int] = {}
    for index, (_, id_lst) in enumerate(sections):
        for entry in id_lst:
            slide_id = entry.get("id")
            if slide_id in live:
                assert slide_id is not None
                owner.setdefault(slide_id, index)

    assignment: dict[str, int] = {}
    current = 0
    for slide_id in order:
        current = owner.get(slide_id, current)
        assignment[slide_id] = current

    for index, (_section, id_lst) in enumerate(sections):
        for entry in list(id_lst):
            id_lst.remove(entry)
        for slide_id in order:
            if assignment[slide_id] == index:
                id_lst.append(
                    parse_xml(f'<p14:sldId xmlns:p14="{_P14_NS}" id="{slide_id}"/>')
                )

    _prune_empty_sections(section_lst)
